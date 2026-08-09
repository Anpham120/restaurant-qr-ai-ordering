# -*- coding: utf-8 -*-
"""Slide bảo vệ môn Học máy và Khai phá dữ liệu.

Dùng lại toàn bộ tiện ích dựng slide của `xuat_slide_pptx.py` — cùng khổ, cùng
bảng màu logo, cùng kiểu bìa — nên bộ slide này nhìn đồng bộ với ba bộ slide
môn khác của dự án.

Ba phần nhóm trưởng yêu cầu, theo đúng thứ tự người nghe cần:

    BÀI TOÁN     hai loại câu hỏi khác nhau về bản chất lời giải
    DỮ LIỆU      thực đơn -> bộ nhãn -> kho tri thức, kèm cách kiểm
    KIẾN TRÚC    bốn đường trả lời và bộ định tuyến

rồi mới tới kết quả đo. Đặt kết quả sau cùng vì mọi con số chỉ có nghĩa khi
người nghe đã biết hệ thống được chia thế nào.

Mỗi slide có kịch bản nói trong phần ghi chú (Presenter View), kèm thời lượng.

Chạy:  python xuat_slide_hoc_may.py
"""
from __future__ import annotations

from pathlib import Path

from pptx.util import Cm, Pt

# Dùng lại nguyên bộ tiện ích, không viết lại: bảng màu, bìa, lộ trình, bảng,
# lưới ảnh, ghi chú nói — tất cả đã có và đã dùng cho ba bộ slide khác.
from xuat_slide_pptx import (  # noqa: E402
    W, H, LE, RONG, DEN, XAM, NHAN, NEN, KE,
    Presentation, PP_ALIGN,
    txt, anh, trang, y_muc, the, bang, chu_thich, luoi_anh,
    bia, lo_trinh, ghi_chu_noi, _luu,
)

HERE = Path(__file__).resolve().parent
BD = HERE / "_bieu_do"                                   # biểu đồ matplotlib
SD = HERE / "output" / "_diagrams_BAO_CAO_HOC_MAY_KP"    # sơ đồ mermaid


def dung() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ══════════════════════════════════════════════════════════ 1 · BÌA
    sl = bia(prs, "ĐỒ ÁN MÔN HỌC", "Học máy và Khai phá dữ liệu",
             "Hệ thống AI tư vấn gọi món cho nhà hàng",
             "Nhóm 05  ·  GVHD: Phạm Ngọc Đông")
    ghi_chu_noi(sl, """
[~40 giây]
Em chào thầy và các bạn. Nhóm em trình bày đồ án môn Học máy và Khai phá dữ
liệu, đề tài Hệ thống AI tư vấn gọi món cho nhà hàng.

Tên đăng ký là "Building a Restaurant Food-Ordering Chatbot using LLM and RAG".
Nhóm giữ nguyên phạm vi đó, và đo thêm được một điều mà đề tài chưa hỏi: RAG nên
dùng ở đâu và KHÔNG nên dùng ở đâu.
""")

    # ══════════════════════════════════════════════════════════ 2 · LỘ TRÌNH
    LT = ["Bài toán đặt ra",
          "Dữ liệu — xây dựng thế nào",
          "Kiến trúc — bốn đường trả lời",
          "Thực nghiệm và kết quả",
          "Hạn chế và kết luận"]
    sl = lo_trinh(prs, LT)
    ghi_chu_noi(sl, """
[~25 giây]
Phần trình bày gồm năm phần. Em bắt đầu từ bài toán, sau đó là cách nhóm xây dữ
liệu, rồi kiến trúc. Kết quả đo để sau cùng, vì mọi con số chỉ có nghĩa khi đã
biết hệ thống được chia thế nào.
""")

    # ══════════════════════════════════════════════ 3 · BÀI TOÁN — bối cảnh
    sl = trang(prs, "Khách quét mã QR tại bàn, và hỏi bằng tiếng Việt tự nhiên",
               "BÀI TOÁN")
    the(sl, [("91", "món trong thực đơn"),
             ("13", "danh mục"),
             ("2", "loại câu hỏi\nkhác nhau về BẢN CHẤT")],
        y=Cm(4.0), cao=Cm(3.6))
    y_muc(sl, [
        (True, "Đủ nhiều để khách không đọc hết thực đơn"),
        (True, "Đủ ít để MỌI câu hỏi đều có đáp án xác định trong dữ liệu"),
        (True, "Nhà hàng phục vụ đồ ăn — một lời khai dị ứng bị bỏ sót "
               "KHÔNG phải lỗi chất lượng, mà là lỗi an toàn"),
    ], y=Cm(8.4), co=18)
    ghi_chu_noi(sl, """
[~50 giây]
Bối cảnh: khách ngồi xuống bàn, quét mã QR, mở giao diện chat.

Thực đơn có 91 món chia 13 danh mục — đủ nhiều để khách không đọc hết, và đủ ít
để mọi câu hỏi đều có đáp án xác định trong dữ liệu.

Điều quan trọng nhất ở đây là dòng cuối. Nhà hàng phục vụ đồ ăn, nên một lời khai
dị ứng bị bỏ sót không phải lỗi chất lượng câu chữ — nó là lỗi an toàn. Ràng buộc
đó định hình toàn bộ kiến trúc mà em trình bày sau.
""")

    # ══════════════════════════════════ 4 · BÀI TOÁN — hai loại câu hỏi
    sl = trang(prs, "Hai loại câu hỏi, khác nhau về BẢN CHẤT lời giải", "BÀI TOÁN")
    bang(sl, ["Loại câu hỏi", "Ví dụ", "Đáp án nằm ở đâu", "Kỹ thuật đúng"], [
        ["Chọn món theo điều kiện", '"Món nào dưới 100 nghìn và không cay?"',
         "Thuộc tính có cấu trúc: giá, nhãn", "Lọc tất định"],
        ["Tri thức nhà hàng", '"Gọi khai vị trước có làm no bụng không?"',
         "Văn xuôi do người viết", "Truy hồi rồi tổng hợp"],
    ], y=Cm(3.9), rong_cot=[0.20, 0.31, 0.26, 0.23], co=15)
    y_muc(sl, [
        (True, "Câu loại một có đáp án ĐÚNG/SAI rõ ràng — `giá < 100.000` là một phép so sánh"),
        (False, "Giao nó cho mô hình sinh là biến bài toán có lời giải chính xác "
                "thành bài toán xấp xỉ"),
        (True, "Câu loại hai KHÔNG có cột nào để lọc — đáp án nằm trong một đoạn văn"),
    ], y=Cm(8.6), co=17)
    ghi_chu_noi(sl, """
[~70 giây]
Khảo sát câu hỏi khách có thể đặt ra, nhóm thấy chúng không đồng nhất về bản chất
lời giải.

Loại thứ nhất là câu chọn món. "Món nào dưới 100 nghìn và không cay" — đáp án nằm
ở thuộc tính có cấu trúc của món. Câu này có đáp án đúng sai rõ ràng, vì "giá nhỏ
hơn 100 nghìn" là một phép so sánh. Giao nó cho mô hình sinh là biến một bài toán
có lời giải chính xác thành một bài toán xấp xỉ.

Loại thứ hai là câu tri thức. "Gọi khai vị trước có làm no bụng không" — không có
cột nào để lọc, đáp án nằm trong một đoạn văn do người viết. Đây mới là chỗ RAG
đúng là câu trả lời.

Phân biệt này không phải chi tiết cài đặt. Nó quyết định toàn bộ kiến trúc.
""")

    # ══════════════════════════════ 5 · BÀI TOÁN — câu hỏi nghiên cứu
    sl = trang(prs, "Câu hỏi nghiên cứu", "BÀI TOÁN")
    o = sl.shapes.add_shape(5, LE, Cm(4.2), RONG, Cm(3.4))
    o.fill.solid(); o.fill.fore_color.rgb = NEN
    o.line.color.rgb = NHAN; o.shadow.inherit = False
    txt(sl, LE + Cm(1.0), Cm(4.9), RONG - Cm(2.0), Cm(2.2),
        "Loại câu hỏi nào KHÔNG nên xử lý bằng RAG,\n"
        "và bằng chứng định lượng nào cho thấy điều đó?",
        24, True, NHAN, PP_ALIGN.CENTER, gian=1.35)
    y_muc(sl, [
        (True, "Câu hỏi KHÔNG phải \"áp dụng RAG cho nhà hàng như thế nào\""),
        (False, "Kỹ thuật RAG đã có sẵn và được dùng rộng rãi"),
        (True, "Đóng góp: phân định RANH GIỚI bằng phép đo, không bằng giả định"),
        (False, "Kèm bốn kết quả âm tính được báo cáo đầy đủ"),
    ], y=Cm(8.6), co=18)
    ghi_chu_noi(sl, """
[~55 giây]
Từ đó nhóm đặt câu hỏi nghiên cứu, và em muốn nói rõ nó KHÔNG phải câu "áp dụng
RAG cho nhà hàng như thế nào". Kỹ thuật RAG đã có sẵn và được dùng rộng rãi; trả
lời câu đó không đóng góp gì mới.

Câu nhóm đặt ra là: loại câu hỏi nào KHÔNG nên xử lý bằng RAG, và bằng chứng định
lượng nào cho thấy điều đó.

Đóng góp của đồ án vì vậy là phân định ranh giới bằng phép đo, kèm bốn kết quả âm
tính — tức bốn lần nhóm thử một thứ nghe rất hợp lý và đo được rằng nó không giúp
gì.
""")

    # ══════════════════════════════════════════ 6 · DỮ LIỆU — tổng quan
    sl = trang(prs, "Ba tầng dữ liệu, dựng theo thứ tự", "DỮ LIỆU")
    the(sl, [("91", "món\n9 trường mỗi món"),
             ("85", "nhãn / 16 họ\n9–21 nhãn mỗi món"),
             ("60", "tài liệu tri thức\n213 đoạn")],
        y=Cm(4.0), cao=Cm(3.6))
    y_muc(sl, [
        (True, "Thực đơn tồn tại ở HAI nơi — tệp JSON cho AI, và CSDL cho backend"),
        (False, "Lúc đầu chúng KHÔNG khớp. Nếu để tự do thì mọi con số của bốn "
                "chặng sau đều đo trên dữ liệu sai"),
        (True, "Nhóm sinh cả hai từ MỘT nguồn, kèm cổng `--check` trong CI"),
        (False, "Điều kiện nghiệm thu: khớp 91/91 món"),
    ], y=Cm(8.4), co=17)
    ghi_chu_noi(sl, """
[~65 giây]
Dữ liệu gồm ba tầng: thực đơn 91 món, bộ nhãn 85 nhãn chia 16 họ, và kho tri thức
60 tài liệu chia 213 đoạn.

Việc đầu tiên phải giải không phải là gán nhãn, mà là một vấn đề đơn giản hơn:
thực đơn tồn tại ở hai nơi — tệp JSON cho dịch vụ AI, và cơ sở dữ liệu cho
backend. Lúc đầu chúng không khớp nhau.

Nếu để hai nguồn tự do thì mọi con số của bốn chặng sau đều đo trên dữ liệu sai.
Nên nhóm sinh cả hai từ một nguồn, kèm một cổng kiểm trong CI, và điều kiện nghiệm
thu là khớp đúng 91 trên 91 món.
""")

    # ══════════════════════════════ 7 · DỮ LIỆU — quy trình gán nhãn
    sl = trang(prs, "Gán nhãn: bốn bước, và ba nguồn có độ tin khác nhau", "DỮ LIỆU")
    bang(sl, ["Nguồn suy ra nhãn", "Nhóm nhãn", "Chủ quan?"], [
        ["Trường có sẵn trong thực đơn", "price, danh mục", "KHÔNG"],
        ["Tên món tự nói ra", "method (Gà NƯỚNG), region (Phở HÀ NỘI)", "rất ít"],
        ["Mô tả món", "flavour, health, allergen", "CÓ"],
    ], y=Cm(3.9), rong_cot=[0.34, 0.44, 0.22], co=15)
    y_muc(sl, [
        (True, "Nhóm thứ ba là nhóm chủ quan — nhóm xử hai cách:"),
        (False, "Bộ rà tự động đối chiếu nhãn với mô tả → tìm ra 7 lỗ thật ở nhóm dị nguyên"),
        (False, "NÓI RA giới hạn trong chính câu trả lời — tài liệu `reading_labels` "
                "viết thẳng rằng nhãn sức khoẻ là đánh giá cảm quan, không phải phân tích"),
    ], y=Cm(9.0), co=17)
    ghi_chu_noi(sl, """
[~75 giây]
Quy trình gán nhãn có bốn bước, nhưng điều đáng nói là ba nguồn suy ra nhãn có độ
tin khác hẳn nhau.

Nhóm thứ nhất lấy từ trường có sẵn — giá, danh mục. Không có gì chủ quan.

Nhóm thứ hai lấy từ tên món. "Gà nướng muối ớt" tự nói ra cách chế biến, "Phở Hà
Nội" tự nói ra vùng miền. Rất ít chủ quan, và kiểm được tự động.

Nhóm thứ ba lấy từ mô tả món — vị, sức khoẻ, dị nguyên. Nhóm này CÓ chủ quan, và
em nói thẳng điều đó.

Nhóm xử hai cách. Thứ nhất là bộ rà tự động đối chiếu nhãn với mô tả, và nó tìm ra
bảy lỗ thật ở nhóm dị nguyên. Thứ hai là nói ra giới hạn trong chính câu trả lời
cho khách.
""")

    # ══════════════════════════ 8 · DỮ LIỆU — nguyên tắc độ phủ
    sl = trang(prs, "Độ phủ quyết định nhãn dùng được vào việc gì", "DỮ LIỆU")
    bang(sl, ["Độ phủ", '"Thiếu nhãn" nghĩa là', "Nhãn dùng để", "Ví dụ họ nhãn"], [
        ["91/91", "LỖI DỮ LIỆU", "LỌC — loại món không thoả",
         "party · meal · season · spice · price"],
        ["một phần", "CHƯA GHI NHẬN, không phải \"không có\"", "SẮP THỨ TỰ — không loại món",
         "occasion · flavour · health · region"],
    ], y=Cm(3.9), rong_cot=[0.11, 0.28, 0.28, 0.33], co=14)
    o = sl.shapes.add_shape(5, LE, Cm(8.4), RONG, Cm(2.6))
    o.fill.solid(); o.fill.fore_color.rgb = NEN
    o.line.color.rgb = NHAN; o.shadow.inherit = False
    txt(sl, LE + Cm(0.8), Cm(8.9), RONG - Cm(1.6), Cm(1.8),
        "Nhãn dị nguyên chỉ phủ 44/91 món\n"
        "→ danh sách lọc ra KHÔNG PHẢI một kết luận về an toàn",
        20, True, NHAN, PP_ALIGN.CENTER, gian=1.3)
    chu_thich(sl, "47 món chưa được ghi nhận dị nguyên nào — không phải là không có. "
                  "Hệ thống nói rõ điều đó với khách thay vì im lặng.")
    ghi_chu_noi(sl, """
[~80 giây]
Đây là nguyên tắc trung tâm của khâu dữ liệu, và nó quyết định nhãn dùng được vào
việc gì.

Nếu một họ nhãn phủ đủ 91 trên 91 món thì thiếu nhãn là lỗi dữ liệu, và nhãn đó
dùng để LỌC — loại món không thoả. Năm họ phủ đủ: nhóm người, bữa ăn, mùa, độ
cay, mức giá.

Nếu một họ chỉ phủ một phần thì thiếu nhãn nghĩa là CHƯA GHI NHẬN, không phải là
KHÔNG CÓ. Nhãn đó chỉ được dùng để sắp thứ tự, không được loại món.

Hệ quả quan trọng nhất nằm ở khung dưới. Nhãn dị nguyên chỉ phủ 44 trên 91 món.
Nghĩa là 47 món chưa được ghi nhận dị nguyên nào — không phải là không có. Nên
danh sách hệ thống lọc ra KHÔNG phải một kết luận về an toàn, và hệ thống nói rõ
điều đó với khách thay vì im lặng.
""")

    # ══════════════════════ 9 · DỮ LIỆU — kho tri thức, hai chế độ
    sl = trang(prs, "Kho tri thức: MỘT kho, HAI chế độ trả lời", "DỮ LIỆU")
    bang(sl, ["Chế độ", "Tài liệu", "Vào chỉ mục", "Cách trả lời", "Mô hình chạm chữ?"], [
        ["verbatim", "24", "KHÔNG", "Tra khoá, trả NGUYÊN VĂN", "0%"],
        ["synthesize", "36", "182 đoạn", "Truy hồi hoặc chọn mục", "không — chỉ trình bày lại"],
    ], y=Cm(3.9), rong_cot=[0.14, 0.12, 0.15, 0.29, 0.30], co=15)
    y_muc(sl, [
        (True, 'Câu "mấy giờ đóng cửa?" có MỘT đáp án đúng duy nhất'),
        (False, "Một chữ số lệch là nói sai sự thật về nhà hàng. Đưa nó qua mô hình "
                "là tạo cơ hội sai ở chỗ chỉ cần đọc ra một chuỗi"),
        (True, "Tài liệu verbatim KHÔNG nằm trong chỉ mục truy hồi"),
        (False, "Nếu để chung thì có HAI đường tới cùng nội dung, và đường xếp hạng "
                "có thể trích câu chính sách ra giữa câu tư vấn món"),
    ], y=Cm(8.6), co=17)
    ghi_chu_noi(sl, """
[~75 giây]
Kho tri thức có 60 tài liệu, nhưng điều đáng nói là nó chia làm hai chế độ trả lời
chứ không chia theo chủ đề.

Hai mươi bốn tài liệu ở chế độ verbatim — trả nguyên văn, mô hình không chạm vào
chữ. Ba mươi sáu tài liệu ở chế độ synthesize — vào chỉ mục truy hồi, 182 đoạn.

Vì sao phải tách? Câu "mấy giờ đóng cửa" có một đáp án đúng duy nhất, và một chữ
số lệch là nói sai sự thật về nhà hàng. Đưa nó qua mô hình là tạo cơ hội sai ở chỗ
chỉ cần đọc ra một chuỗi.

Hệ quả kiến trúc: tài liệu verbatim KHÔNG nằm trong chỉ mục truy hồi. Nếu để chung
thì có hai đường tới cùng một nội dung, và đường xếp hạng có thể trích một câu
chính sách ra giữa câu tư vấn món. Có test chốt điều này.
""")

    # ══════════════════ 10 · DỮ LIỆU — tám tài liệu do máy sinh
    sl = trang(prs, "Tám tài liệu chứa con số do MÁY SINH, không viết tay", "DỮ LIỆU")
    y_muc(sl, [
        (True, "Văn xuôi kể lại dữ liệu thì LUÔN trôi khỏi dữ liệu"),
        (False, 'Một tài liệu viết tay ghi "hơn 90 món" trong khi thực đơn có đúng 91 '
                '— sai ngay từ lúc viết, và không ai canh'),
        (True, "Tám tài liệu có số được sinh lại từ thực đơn mỗi lần, kèm cổng `--check`"),
        (False, "menu_size · price_range · preorder · takeaway_items · children · "
                "vegetarian · spice_levels · allergen_labelling"),
        (True, "Nhờ vậy chúng KHÔNG THỂ lệch — thực đơn đổi thì tài liệu tự đổi theo"),
    ], y=Cm(3.9), co=18)
    chu_thich(sl, "Mười sáu tài liệu verbatim còn lại là chính sách thật của nhà hàng, "
                  "không suy được từ thực đơn nên phải viết tay.")
    ghi_chu_noi(sl, """
[~65 giây]
Trong 24 tài liệu verbatim, có tám tài liệu chứa con số, và tám tài liệu đó do máy
sinh chứ không viết tay.

Lý do là một bài học nhóm đã trả giá: văn xuôi kể lại dữ liệu thì luôn trôi khỏi
dữ liệu. Một tài liệu viết tay của bản đầu ghi "hơn 90 món" trong khi thực đơn có
đúng 91. Sai ngay từ lúc viết, và không ai canh.

Tám tài liệu này — số món, dải giá, số món cho trẻ em, số món chay — được sinh lại
từ thực đơn mỗi lần, kèm một cổng kiểm trong CI. Nhờ vậy chúng không thể lệch:
thực đơn thêm một món thì tài liệu tự cập nhật số đếm.

Mười sáu tài liệu còn lại là chính sách thật của nhà hàng, không suy được từ thực
đơn nên phải viết tay.
""")

    # ═══════════════════════ 11 · KIẾN TRÚC — bốn đường trả lời
    sl = trang(prs, "Bốn đường trả lời, phân theo mức được phép tin mô hình",
               "KIẾN TRÚC")
    if (SD / "so-do-2.png").exists():
        anh(sl, SD / "so-do-2.png", LE, Cm(3.7), RONG * 0.52, Cm(9.6))
    bang(sl, ["Đường", "Rủi ro chệch"], [
        ["Lọc nhãn — không đọc kho", "KHÔNG CÓ"],
        ["Tra khoá — trả nguyên văn", "KHÔNG CÓ"],
        ["Chọn mục — trong 1 tài liệu", "thấp"],
        ["Truy hồi — toàn kho 182 đoạn", "cao nhất"],
    ], y=Cm(4.6), rong_cot=[0.66, 0.34], co=14)
    # đẩy bảng sang nửa phải
    for sh in sl.shapes:
        if sh.has_table:
            sh.left = int(LE + RONG * 0.55)
            sh.width = int(RONG * 0.45)
    chu_thich(sl, "Đường càng hay dùng thì càng ít rủi ro — đó là chủ ý, "
                  "không phải tình cờ.")
    ghi_chu_noi(sl, """
[~70 giây]
Đây là kiến trúc trung tâm. Hệ thống có bốn đường trả lời, và chúng được phân theo
mức được phép tin mô hình.

Đường thứ nhất là lọc nhãn — đọc thẳng thuộc tính của thực đơn, không chạm kho tri
thức. Không có rủi ro chệch nào, vì nó tra bảng.

Đường thứ hai là tra khoá — trả nguyên văn tài liệu chính sách. Cũng không có rủi
ro, vì không có xếp hạng.

Đường thứ ba là chọn mục trong một tài liệu đã biết — có xếp hạng, nhưng phạm vi
chỉ ba tới tám đoạn.

Đường thứ tư là truy hồi toàn kho 182 đoạn — rủi ro cao nhất.

Điều đáng nói ở dòng cuối: cách sắp này có chủ ý. Đường càng hay dùng thì càng ít
rủi ro.
""")

    # ═══════════════════════ 12 · KIẾN TRÚC — bộ định tuyến
    sl = trang(prs, "Bộ định tuyến: chuỗi cổng có thứ tự, KHÔNG phải bộ phân loại",
               "KIẾN TRÚC")
    if (SD / "so-do-5.png").exists():
        anh(sl, SD / "so-do-5.png", LE, Cm(3.7), RONG * 0.34, Cm(9.8))
    y_muc(sl, [
        (True, "Không mô hình, không điểm tin cậy, không argmax"),
        (False, "Cổng nào khớp trước thì thắng — mỗi vị trí đứng ở đó vì một ca hỏng đo được"),
        (True, "Nhánh xã giao PHẢI đứng trước mọi nhánh chọn món"),
        (False, 'Thiếu nó thì "xin chào" rơi xuống truy hồi — vì sau khi rút dấu, '
                '"chao" khớp món "Cháo lòng Sài Gòn"'),
        (True, "Truy hồi đứng GẦN CUỐI — RAG là phương án cuối, không phải mặc định"),
    ], y=Cm(4.0), x=int(LE + RONG * 0.38), w=int(RONG * 0.62), co=16)
    ghi_chu_noi(sl, """
[~75 giây]
Bộ định tuyến quyết định câu hỏi đi đường nào, và điều đầu tiên em muốn nói là nó
KHÔNG phải một bộ phân loại. Không có mô hình, không có điểm tin cậy, không có
argmax.

Nó là một chuỗi cổng có thứ tự cố định, cổng nào khớp trước thì thắng. Và mỗi vị
trí trong chuỗi đứng ở đó vì một ca hỏng đo được.

Ví dụ rõ nhất: nhánh xã giao phải đứng trước mọi nhánh chọn món. Thiếu nó thì câu
"xin chào" rơi xuống truy hồi, và khách nhận về một danh sách rượu nếp cẩm. Lý do
là cổng kiểm thuộc miền so từng từ đơn sau khi rút dấu, nên chữ "chao" khớp với
món "Cháo lòng Sài Gòn".

Và truy hồi đứng gần cuối. Đó là chủ ý: RAG là phương án cuối, không phải phương
án mặc định.
""")

    # ═══════════════════ 13 · KIẾN TRÚC — hiểu câu hỏi
    sl = trang(prs, "Hiểu câu hỏi: 629 cụm từ vựng, không dùng mô hình", "KIẾN TRÚC")
    if (SD / "so-do-7.png").exists():
        anh(sl, SD / "so-do-7.png", LE, Cm(3.7), RONG * 0.30, Cm(9.6))
    y_muc(sl, [
        (True, "Rút dấu tiếng Việt là phép MẤT thông tin"),
        (False, '"bò" và "bơ" cùng thành "bo" — 107 trong 629 cụm có nguy cơ đụng chữ'),
        (True, "Cách chặn là MỘT LUẬT, không phải danh sách ngoại lệ"),
        (False, "Khớp cụm DÀI trước, rồi ĂN HẾT đoạn đã khớp"),
        (True, "Ràng buộc khác ngữ cảnh — và nhầm hai thứ này gây lỗi an toàn"),
        (False, '"không cay" LOẠI món · "đi hẹn hò" chỉ XẾP LÊN TRƯỚC'),
    ], y=Cm(4.0), x=int(LE + RONG * 0.34), w=int(RONG * 0.66), co=16)
    ghi_chu_noi(sl, """
[~80 giây]
Lớp hiểu câu hỏi biến chuỗi tiếng Việt thành một cấu trúc có 48 trường. Nó dùng
629 cụm từ vựng và KHÔNG dùng mô hình.

Vấn đề lớn nhất của khâu này là rút dấu tiếng Việt. Người Việt gõ không dấu rất
thường, nên phải rút dấu để khớp được. Nhưng rút dấu là phép mất thông tin: "bò"
và "bơ" cùng thành "bo".

Nhóm kiểm kê được 107 trong 629 cụm có nguy cơ đụng chữ. Cách chặn không phải vá
từng lỗi, mà là một luật: khớp cụm dài trước, rồi ăn hết đoạn đã khớp. Nhờ vậy
thêm cụm "mì chính" là tự động chặn được cụm "mì" khớp nhầm vào giữa.

Chỗ khó thứ hai là phân biệt ràng buộc với ngữ cảnh. "Không cay" là ràng buộc —
món cay phải bị loại. "Đi hẹn hò" là ngữ cảnh — món hợp dịp chỉ được xếp lên
trước. Nhầm hai thứ này thì hoặc lọc mất món đúng, hoặc để lọt món khách không ăn
được.
""")

    # ═══════════════════ 14 · KIẾN TRÚC — bốn lớp an toàn
    sl = trang(prs, "An toàn bảo đảm bằng CẤU TRÚC, không bằng lời nhắc mô hình",
               "KIẾN TRÚC")
    the(sl, [("1", "Lọc dị nguyên\nFAIL-CLOSED\náp cuối, không bao giờ nới"),
             ("2", "Danh sách trắng nhánh\n2/19 nhánh được sinh chữ"),
             ("3", "10 phép xác minh\nvi phạm thì BỎ CẢ CÂU"),
             ("4", "Thẻ giỏ TẤT ĐỊNH\nkhông đọc chữ mô hình viết")],
        y=Cm(4.0), cao=Cm(4.2), co_so=34)
    y_muc(sl, [
        (True, "Bật đường sinh làm 15 ca tụt — 14 trong 15 là ca DỊ NGUYÊN"),
        (False, 'Mô hình viết văn mượt hơn và BỎ câu "bạn nhắc nhân viên để bếp xác nhận". '
                'Câu đó là NỘI DUNG, không phải văn vẻ — vì nhãn dị nguyên chỉ phủ 44/91'),
        (True, "Lời nhắc là ĐỀ NGHỊ. Chỉ phép kiểm sau khi sinh mới là BẢO ĐẢM"),
    ], y=Cm(9.0), co=16)
    ghi_chu_noi(sl, """
[~85 giây]
An toàn được bảo đảm bằng bốn lớp độc lập, và điểm chung của cả bốn là chúng nằm
trong cấu trúc chứ không nằm trong lời nhắc mô hình.

Lớp một: lọc dị nguyên fail-closed, áp cuối cùng và không bao giờ nới, kể cả khi
kết quả rỗng. Thà nói không có món nào phù hợp còn hơn mời một món có thể gây hại.

Lớp hai: danh sách trắng nhánh. Chỉ hai trong mười chín nhánh được phép sinh chữ.
Mười bảy nhánh còn lại không có đường để mô hình ghi chữ cho khách.

Lớp ba: mười phép xác minh. Vi phạm bất kỳ phép nào thì bỏ cả câu sinh, dùng lại
câu khuôn mẫu — không sửa, không thử lại.

Lớp bốn: thẻ giỏ hàng dựng từ danh sách món đã lọc, không đọc chữ mô hình viết.

Dòng dưới là bằng chứng cho thấy vì sao phải làm vậy. Khi bật đường sinh lần đầu,
15 ca tụt, và 14 trong 15 là ca dị ứng. Lý do: mô hình viết văn mượt hơn nên nó bỏ
câu "bạn nhắc nhân viên để bếp xác nhận". Câu đó là nội dung chứ không phải văn
vẻ, vì nhãn dị nguyên chỉ phủ 44 trên 91 món.

Lời nhắc cũng đã yêu cầu điều đó. Nhưng yêu cầu trong lời nhắc là đề nghị, không
phải bảo đảm.
""")

    # ═══════════════ 15 · KẾT QUẢ — so ba phương pháp truy hồi
    sl = trang(prs, "So ba phương pháp truy hồi — 66 ca văn xuôi viết tay",
               "KẾT QUẢ")
    if (BD / "bd1-truy-hoi.png").exists():
        anh(sl, BD / "bd1-truy-hoi.png", LE, Cm(3.6), RONG * 0.60, Cm(10.2))
    y_muc(sl, [
        (True, "Chốt embedding bge-m3"),
        (False, "Hybrid cao hơn ở Hit@1 nhưng THUA ở Hit@2 — cột hệ thống thật sự dùng"),
        (True, "Chấm ở k khác k hệ thống dùng là đo một hệ thống KHÔNG TỒN TẠI"),
        (True, "cấm@5 quan trọng hơn Hit@5"),
        (False, "Nó đo việc trả lời SAI, không phải kém"),
    ], y=Cm(4.4), x=int(LE + RONG * 0.63), w=int(RONG * 0.37), co=15)
    ghi_chu_noi(sl, """
[~80 giây]
Sang phần kết quả. Đây là phép so ba phương pháp truy hồi trên 66 ca nhắm vào văn
xuôi viết tay — bài toán RAG thật của hệ thống.

Nhóm chốt embedding bge-m3. Nhưng điều đáng nói nằm ở cột được tô đậm.

Nếu nhìn Hit@1 thì hybrid cao hơn — 0,712 so với 0,697. Nhưng hệ thống lúc chạy
trích đúng HAI đoạn, nên cột quyết định là Hit@2. Ở cột đó embedding đạt 0,879 còn
hybrid chỉ 0,803.

Nếu nhóm chốt theo Hit@1 thì đã chọn phương pháp kém hơn cho chính hệ thống của
mình, và bảng số vẫn trông đúng. Chấm ở k khác k hệ thống dùng là đo một hệ thống
không tồn tại.

Cột cuối là cấm@5 — số ca lấy phải đoạn thuộc chủ đề mà câu hỏi không được chạm.
Nó đo việc trả lời SAI, không phải kém, nên nhóm đặt nó cao hơn Hit@5.
""")

    # ═══════════════ 16 · KẾT QUẢ — RAG chạy bao nhiêu
    sl = trang(prs, "Kết quả trung tâm: RAG chạy bao nhiêu trong một phiên thật?",
               "KẾT QUẢ")
    if (BD / "bd4-duong-di.png").exists():
        anh(sl, BD / "bd4-duong-di.png", LE, Cm(3.6), RONG * 0.62, Cm(10.4))
    y_muc(sl, [
        (True, "96,9% lượt KHÔNG cần chạm kho tri thức"),
        (False, "Chúng là câu chọn món, và lọc tất định trả lời chính xác 100,00%"),
        (True, "Không có nghĩa RAG vô dụng"),
        (False, "Nghĩa là hai tập đó được viết quanh các nhánh tất định"),
        (True, "Nhóm đã bổ sung tập ca chạy RAG trong phiên có bộ nhớ"),
        (False, "12/12 lượt đạt, 0 lỗi an toàn"),
    ], y=Cm(4.2), x=int(LE + RONG * 0.65), w=int(RONG * 0.35), co=15)
    ghi_chu_noi(sl, """
[~85 giây]
Đây là kết quả trung tâm của đồ án, và nó là phép đo đổi cách đọc mọi con số phía
trước.

Chạy 163 lượt kịch bản như một phiên thật có bộ nhớ: 96,9% lượt KHÔNG cần chạm tới
kho tri thức. Chúng là câu chọn món, và phép lọc tất định trả lời chúng chính xác
100 phần trăm.

Cột đáng nhìn nhất là cột bằng 0 — truy hồi toàn kho.

Nhưng em phải nói rõ điều này KHÔNG có nghĩa RAG vô dụng. Nó có nghĩa hai tập đánh
giá đó được viết quanh các nhánh tất định, nên chúng không hỏi tới nhánh truy hồi.

Và đó chính là một lỗ của tập ca chứ không phải của hệ thống. Nhóm đã bổ sung một
nhóm kịch bản đặt câu tri thức ở giữa phiên, sau lời khai dị ứng. Kết quả 12 trên
12 lượt đạt, 0 lỗi an toàn — nên giờ có bằng chứng nhánh tri thức chạy được, thay
vì chỉ có bằng chứng nó không chạy.
""")

    # ═══════════════ 17 · KẾT QUẢ — bốn kết quả âm tính
    sl = trang(prs, "Bốn kết quả âm tính — báo cáo đầy đủ, không giấu", "KẾT QUẢ")
    bang(sl, ["Đã thử", "McNemar", "Kết luận"], [
        ["Hybrid BM25 + embedding", "p = 1,0000", "Hoà tuyệt đối — 18 ca sửa được, 18 ca làm hỏng"],
        ["Xếp hạng lại cross-encoder", "p = 0,8238", "Hoà, và CHẬM 118 lần"],
        ["Gộp tài liệu sinh theo nhãn", "p = 0,5488", "Hoà — không đổi cấu trúc"],
        ["Bỏ hẳn nhóm tài liệu đó", "—", "BỎ ĐƯỢC sau khi từ vựng đưa 99,1% câu về nhánh lọc"],
    ], y=Cm(3.9), rong_cot=[0.28, 0.16, 0.56], co=14)
    y_muc(sl, [
        (True, "Ba cách chữa ĐỘC LẬP đều không thắng"),
        (False, "Khi ba cách chữa độc lập cùng hoà, vấn đề không nằm ở cách chữa "
                "mà ở CHẨN ĐOÁN"),
        (True, "Kiến trúc cuối GỌN HƠN kiến trúc đầu — và đó là kết quả của đo lường"),
    ], y=Cm(9.4), co=17)
    ghi_chu_noi(sl, """
[~80 giây]
Nhóm báo cáo đầy đủ bốn kết quả âm tính, tức bốn lần thử một thứ nghe rất hợp lý
và đo được rằng nó không giúp gì.

Thứ nhất: hybrid trộn BM25 với embedding. p bằng 1,0000 — hoà tuyệt đối. 18 ca sửa
được, 18 ca làm hỏng. Khó có kết quả nào sạch hơn thế.

Thứ hai: xếp hạng lại bằng cross-encoder, công cụ chuẩn cho đúng vấn đề nhóm chẩn
đoán. p bằng 0,8238, và chậm 118 lần. Ngay cả khi nó thắng thì con số đó cũng
không triển khai được.

Thứ ba: gộp các tài liệu sinh theo nhãn lại. p bằng 0,5488.

Ba cách chữa độc lập đều hoà. Và khi ba cách chữa độc lập cùng hoà thì vấn đề
không nằm ở cách chữa mà ở chẩn đoán. Nhóm bỏ hẳn nhóm tài liệu đó, và Hit@2 lên
0,879.

Kết quả là kiến trúc cuối GỌN HƠN kiến trúc đầu — và đó là kết quả của đo lường,
không phải của việc cắt bớt cho kịp.
""")

    # ═══════════════ 18 · KẾT QUẢ — bảng tổng hợp
    sl = trang(prs, "Tổng hợp kết quả — mọi tệp bằng chứng cùng ngày 09/08/2026",
               "KẾT QUẢ")
    bang(sl, ["Phép đo", "Quy mô", "Kết quả"], [
        ["Tập ca trả lời", "147 ca", "147/147 · niêm phong 48/48"],
        ["Bộ nhớ phiên", "63 kịch bản / 175 lượt", "không lượt nào đỏ · 0 lỗi an toàn"],
        ["Golden đầu-cuối qua STACK THẬT", "103 lượt × 2 cấu hình", "103/103 ở CẢ HAI"],
        ["Truy hồi nhóm written", "66 ca", "embedding Hit@2 0,879 · cấm@5 6"],
        ["LLM + RAG câu loại C", "76 ca", "76/76 cả hai chiều · 0 ca tụt"],
        ["Chọn món", "50 câu", "lọc nhãn 100,00% · 0 món vi phạm"],
        ["Bộ kiểm", "—", "429 + 143 test · 14 cổng --check"],
    ], y=Cm(3.7), rong_cot=[0.36, 0.26, 0.38], co=13)
    chu_thich(sl, "Golden đo trên stack dựng lại TỪ SỐ KHÔNG: ảnh AI build mới, "
                  "cơ sở dữ liệu migrate từ trống rỗng.")
    ghi_chu_noi(sl, """
[~70 giây]
Đây là bảng tổng hợp, và em muốn nhấn một điểm: mọi tệp bằng chứng đều cùng ngày,
mùng 9 tháng 8.

147 trên 147 ca trả lời, trong đó tập niêm phong 48 trên 48.

175 lượt bộ nhớ phiên, không lượt nào đỏ, 0 lỗi an toàn.

Dòng thứ ba là dòng quan trọng nhất: golden 103 trên 103 ở CẢ HAI cấu hình, đo qua
stack thật — quét QR, backend, dịch vụ AI, thẻ giỏ, giỏ hàng thật. Và như ghi chú
dưới, stack đó được dựng lại từ số không: ảnh AI build mới, cơ sở dữ liệu migrate
từ trống rỗng.

Việc dựng lại từ số không tìm ra ba lỗi triển khai mà bản dựng cũ che mất — không
lỗi nào là lỗi hệ thống, cả ba là lỗi quy trình dựng.
""")

    # ═══════════════════════════ 19 · HẠN CHẾ
    sl = trang(prs, "Hạn chế — nói ra thay vì giấu", "KẾT LUẬN")
    y_muc(sl, [
        (True, "KHÔNG có nhật ký hội thoại của khách thật"),
        (False, "Mọi ca đánh giá do nhóm viết. Số đo được hệ thống CÓ TÔN TRỌNG RÀNG BUỘC "
                "hay không; nó KHÔNG đo được khách thật sẽ hỏi gì"),
        (True, "Cả bốn tập niêm phong đã mở"),
        (False, "Không con số nào còn là held-out. Muốn kết luận held-out lần nữa "
                "thì phải viết tập MỚI"),
        (True, "Nhãn dị nguyên phủ 44/91, và mô tả món không phải bảng thành phần"),
        (False, "Còn thiếu bao nhiêu thì KHÔNG BIẾT ĐƯỢC từ dữ liệu này"),
        (True, "Chỉ hiểu tiếng Việt — và giới hạn này chạm tới an toàn"),
        (False, "Lời khai dị ứng bằng tiếng Anh KHÔNG bật hàng rào dị nguyên"),
    ], y=Cm(3.8), co=16)
    ghi_chu_noi(sl, """
[~85 giây]
Phần hạn chế, và nhóm nói ra thay vì giấu.

Hạn chế lớn nhất: không có nhật ký hội thoại của khách thật. Mọi ca đánh giá đều do
nhóm viết. Nghĩa là con số đo được hệ thống có tôn trọng ràng buộc hay không, nhưng
nó không đo được khách thật sẽ hỏi những gì. Đây là hạn chế không sửa được bằng
cách viết thêm ca.

Thứ hai: cả bốn tập niêm phong đã được mở trong quá trình làm. Nên không con số nào
trong báo cáo còn là held-out.

Thứ ba: nhãn dị nguyên phủ 44 trên 91 món, và mô tả món không phải bảng thành phần.
Bộ rà tìm được chỗ mô tả có nhắc mà nhãn thiếu; nó không tìm được món có dị nguyên
mà mô tả cũng không nhắc.

Thứ tư, và đây là hạn chế chạm tới an toàn: hệ thống chỉ hiểu tiếng Việt. Một lời
khai dị ứng bằng tiếng Anh KHÔNG bật hàng rào dị nguyên, trong khi câu tiếng Việt
tương đương thì bật.
""")

    # ═══════════════════════════ 20 · KẾT LUẬN
    sl = trang(prs, "Kết luận", "KẾT LUẬN")
    o = sl.shapes.add_shape(5, LE, Cm(3.9), RONG, Cm(3.0))
    o.fill.solid(); o.fill.fore_color.rgb = NEN
    o.line.color.rgb = NHAN; o.shadow.inherit = False
    txt(sl, LE + Cm(1.0), Cm(4.5), RONG - Cm(2.0), Cm(2.0),
        "Kiến trúc cuối GỌN HƠN kiến trúc đầu,\n"
        "và đó là kết quả của ĐO LƯỜNG chứ không phải của việc cắt bớt cho kịp.",
        21, True, NHAN, PP_ALIGN.CENTER, gian=1.35)
    y_muc(sl, [
        (True, "Một mô hình nhúng thay vì ba phương pháp truy hồi — hybrid p = 1,0000"),
        (True, "Không xếp hạng lại — p = 0,8238 và chậm 118 lần"),
        (True, "Kho tri thức nhỏ đi 3/5 — ba cách chữa độc lập đều hoà"),
        (True, "Bài học phương pháp đáng giữ nhất:"),
        (False, "Trong đồ án này, số lần THƯỚC ĐO sai nhiều hơn số lần HỆ THỐNG sai. "
                "Thứ phát hiện ra chúng là kỷ luật in dữ liệu thô kèm tỷ lệ — "
                "một bảng số không tự tố cáo được nó sai, còn vài dòng ví dụ thì có"),
    ], y=Cm(7.6), co=16)
    ghi_chu_noi(sl, """
[~75 giây]
Để kết luận.

Kết quả chính của đồ án là một hệ thống ĐƠN GIẢN HƠN thiết kế ban đầu: một mô hình
nhúng thay vì ba phương pháp truy hồi, không xếp hạng lại, kho tri thức nhỏ đi ba
phần năm.

Và mỗi lần bỏ bớt đều có một phép đo nói rằng thêm cái này không giúp gì — hybrid
p bằng 1,0000, reranker p bằng 0,8238, gộp tài liệu p bằng 0,5488.

Nên câu kết là: kiến trúc cuối gọn hơn kiến trúc đầu, và đó là kết quả của đo lường
chứ không phải của việc cắt bớt cho kịp.

Bài học phương pháp đáng giữ nhất: trong đồ án này, số lần thước đo sai nhiều hơn
số lần hệ thống sai. Thứ phát hiện ra chúng không phải sự cẩn thận, mà là kỷ luật
in dữ liệu thô kèm tỷ lệ. Một bảng số không tự tố cáo được nó sai, còn vài dòng ví
dụ thì có.

Em xin hết. Nhóm em sẵn sàng nhận câu hỏi của thầy.
""")

    ra = HERE / "output" / "SLIDE_HOC_MAY_KPDL.pptx"
    return _luu(prs, ra)


if __name__ == "__main__":
    print(dung())
