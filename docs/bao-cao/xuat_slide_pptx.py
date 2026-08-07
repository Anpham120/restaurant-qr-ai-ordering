"""Dựng slide giới thiệu sản phẩm CMC Restaurant (.pptx, khổ 16:9).

Mạch trình bày đi theo sản phẩm chứ không theo báo cáo:
sản phẩm là gì → ai dùng → chức năng theo vai trò → kiến trúc và công nghệ
→ điểm kỹ thuật đáng nói → chất lượng và vận hành → hạn chế → hướng phát triển.

Ảnh và sơ đồ dùng lại đúng tệp của báo cáo nên slide không thể lệch nội dung.
Chạy ``render_so_do.py`` trước nếu chưa có sơ đồ.

    python docs/bao-cao/xuat_slide_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

HERE = Path(__file__).resolve().parent
RA = HERE / "output" / "SLIDE_GIOI_THIEU_SAN_PHAM.pptx"
ANH = HERE / "../assets/report"
RM = HERE / "../assets/readme"
SO_DO = HERE / "output" / "_diagrams"
LOGO = HERE / "../../frontend/src/mocks/images/logo.png"

W, H = Cm(33.867), Cm(19.05)
LE = Cm(1.6)
RONG = W - 2 * LE

# Bảng màu lấy theo logo nhà hàng: nâu đậm và vàng đồng.
DEN = RGBColor(0x11, 0x11, 0x11)
XAM = RGBColor(0x5A, 0x5A, 0x5A)
NHAN = RGBColor(0x6B, 0x3E, 0x2E)          # nâu logo
VANG = RGBColor(0xB0, 0x86, 0x2A)          # vàng đồng
NEN = RGBColor(0xF7, 0xF4, 0xF1)
NEN_DAM = RGBColor(0xEE, 0xE4, 0xDC)
KE = RGBColor(0xDC, 0xD2, 0xC9)
FONT = "Times New Roman"


# ------------------------------------------------------------------ tiện ích

def txt(sl, x, y, w, h, chu, co=17, dam=False, mau=DEN,
        can=PP_ALIGN.LEFT, gian=1.2, nghieng=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for k, dong in enumerate(chu.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = can
        p.line_spacing = gian
        r = p.add_run()
        r.text = dong
        r.font.size = Pt(co)
        r.font.bold = dam
        r.font.italic = nghieng
        r.font.color.rgb = mau
        r.font.name = FONT
    return tb


def anh(sl, f, x, y, w_max, h_max, vien=True):
    f = Path(f)
    with Image.open(f) as im:
        pw, ph = im.size
    w, h = w_max, w_max * ph / pw
    if h > h_max:
        h, w = h_max, h_max * pw / ph
    pic = sl.shapes.add_picture(str(f), int(x + (w_max - w) // 2),
                                int(y + (h_max - h) // 2), width=int(w))
    if vien:
        pic.line.color.rgb = KE
        pic.line.width = Pt(0.75)
    return pic


def dau_logo(sl, cao=Cm(1.55)):
    """Logo nhà hàng đặt nhỏ ở góc phải trên mỗi slide nội dung."""
    if not LOGO.exists():
        return
    with Image.open(LOGO) as im:
        pw, ph = im.size
    w = cao * pw / ph
    sl.shapes.add_picture(str(LOGO), int(W - LE - w), Cm(0.75), height=int(cao))


def trang(prs, tieu_de, nhan=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rong_chu = RONG - Cm(2.4)          # chừa chỗ cho logo góc phải
    if nhan:
        txt(sl, LE, Cm(0.95), rong_chu, Cm(0.8), nhan, 13, True, NHAN)
        txt(sl, LE, Cm(1.6), rong_chu, Cm(1.3), tieu_de, 26, True, DEN)
    else:
        txt(sl, LE, Cm(1.2), rong_chu, Cm(1.5), tieu_de, 28, True, DEN)
    v = sl.shapes.add_shape(1, LE, Cm(3.0), RONG, Cm(0.045))
    v.fill.solid(); v.fill.fore_color.rgb = NHAN
    v.line.fill.background(); v.shadow.inherit = False
    dau_logo(sl)
    return sl


def y_muc(sl, muc, y=Cm(3.8), x=LE, w=None, co=17, gian=1.45):
    w = w or RONG
    tb = sl.shapes.add_textbox(x, y, w, H - y - Cm(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for k, (chinh, chu) in enumerate(muc):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.line_spacing = gian
        p.space_after = Pt(7 if chinh else 10)
        if chinh:
            r = p.add_run(); r.text = "▪  "
            r.font.size = Pt(co); r.font.color.rgb = NHAN; r.font.name = FONT
        else:
            p.level = 1
        r = p.add_run(); r.text = chu
        r.font.size = Pt(co if chinh else co - 2)
        r.font.color.rgb = DEN if chinh else XAM
        r.font.name = FONT
    return tb


def the(sl, muc, y=Cm(3.9), cao=Cm(3.2), co_so=28):
    n = len(muc)
    khe = Cm(0.45)
    w = (RONG - khe * (n - 1)) // n
    for i, (so, nh) in enumerate(muc):
        o = sl.shapes.add_shape(5, LE + i * (w + khe), y, w, cao)
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        tf = o.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = so
        r.font.size = Pt(co_so); r.font.bold = True
        r.font.color.rgb = DEN; r.font.name = FONT
        for d in nh.split("\n"):
            q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
            rr = q.add_run(); rr.text = d
            rr.font.size = Pt(12); rr.font.color.rgb = XAM; rr.font.name = FONT


def bang(sl, tieu_de, hang, y=Cm(3.8), rong_cot=None, co=14):
    n = len(tieu_de)
    t = sl.shapes.add_table(len(hang) + 1, n, LE, y, RONG,
                            Cm(1.0) * (len(hang) + 1)).table
    if rong_cot:
        for i, c in enumerate(rong_cot):
            t.columns[i].width = int(RONG * c)
    for i, x in enumerate(tieu_de):
        o = t.cell(0, i)
        o.text = x
        o.fill.solid(); o.fill.fore_color.rgb = NEN_DAM
        for p in o.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(co); r.font.bold = True
                r.font.color.rgb = DEN; r.font.name = FONT
    for j, h in enumerate(hang, 1):
        for i, x in enumerate(h):
            o = t.cell(j, i)
            o.text = x
            o.fill.solid()
            o.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in o.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(co - 1)
                    r.font.bold = (i == 0)
                    r.font.color.rgb = DEN if i == 0 else XAM
                    r.font.name = FONT
    return t


def chu_thich(sl, chu):
    txt(sl, LE, H - Cm(1.3), RONG, Cm(0.9), chu, 12, False, XAM,
        nghieng=True)


def luoi_anh(sl, items, y, cao, x0=None, w0=None):
    """items: [(đường_ảnh, nhãn)] — xếp thành một hàng ngang.

    Nhãn rỗng thì không tạo hộp chữ, tránh sinh hộp trống tràn khỏi khung.
    """
    x0 = LE if x0 is None else x0
    w0 = RONG if w0 is None else w0
    n = len(items)
    khe = Cm(0.5)
    w = (w0 - khe * (n - 1)) // n
    for i, (f, nh) in enumerate(items):
        x = x0 + i * (w + khe)
        anh(sl, f, x, y, w, cao)
        if nh:
            txt(sl, x, y + cao + Cm(0.15), w, Cm(1.0), nh, 13, True, DEN,
                PP_ALIGN.CENTER, gian=1.1)





# ------------------------------------------------------------------ nội dung

SD_DA = HERE / "output" / "_diagrams_BAO_CAO_DO_AN_CHUY"     # sơ đồ đồ án
SD_WEB = HERE / "output" / "_diagrams_BAO_CAO_CONG_NGHE_"    # sơ đồ môn web


def ghi_chu_noi(sl, chu: str) -> None:
    """Kịch bản nói, đặt vào phần ghi chú của slide (Presenter View)."""
    sl.notes_slide.notes_text_frame.text = chu.strip()


def bia(prs, loai, mon, ten_de_tai, phu=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    txt(sl, LE, Cm(1.1), RONG, Cm(1), "TRƯỜNG ĐẠI HỌC CMC · KHOA CNTT & TT",
        15, True, NHAN, PP_ALIGN.CENTER)
    txt(sl, LE, Cm(2.1), RONG, Cm(1), mon, 14, False, XAM, PP_ALIGN.CENTER)
    if LOGO.exists():
        with Image.open(LOGO) as im:
            pw, ph = im.size
        cao = Cm(4.4)
        sl.shapes.add_picture(str(LOGO), int((W - cao * pw / ph) // 2), Cm(3.3),
                              height=int(cao))
    txt(sl, LE, Cm(8.2), RONG, Cm(1.2), loai, 17, True, NHAN, PP_ALIGN.CENTER)
    txt(sl, LE, Cm(9.4), RONG, Cm(2.4), ten_de_tai, 24, True, DEN,
        PP_ALIGN.CENTER, gian=1.3)
    if phu:
        txt(sl, LE, Cm(12.6), RONG, Cm(1), phu, 14, False, NHAN, PP_ALIGN.CENTER)
    txt(sl, LE, Cm(14.4), RONG, Cm(1.2),
        "Nhóm 05 sinh viên   ·   GVHD: Ngô Việt Anh", 15, False, XAM,
        PP_ALIGN.CENTER)
    txt(sl, LE, Cm(16.8), RONG, Cm(1), "Hà Nội, tháng 8 năm 2026", 13, False,
        XAM, PP_ALIGN.CENTER)
    return sl


def lo_trinh(prs, muc, danh_dau=None):
    """Slide bản đồ nội dung, tô đậm phần sắp trình bày."""
    sl = trang(prs, "Nội dung trình bày", "LỘ TRÌNH")
    khe = Cm(0.4)
    n = len(muc)
    cao = (H - Cm(4.6)) / n
    for i, m in enumerate(muc):
        y = Cm(3.7) + i * (cao)
        noi_bat = (danh_dau is not None and i == danh_dau)
        o = sl.shapes.add_shape(5, LE, y, RONG, cao - khe)
        o.fill.solid()
        o.fill.fore_color.rgb = NEN_DAM if noi_bat else NEN
        o.line.color.rgb = NHAN if noi_bat else KE
        o.shadow.inherit = False
        tf = o.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"   {i + 1}.   {m}"
        r.font.size = Pt(16)
        r.font.bold = noi_bat
        r.font.color.rgb = NHAN if noi_bat else DEN
        r.font.name = FONT
    return sl


# ==========================================================================
# BỘ 1 — ĐỒ ÁN CHUYÊN NGÀNH
# ==========================================================================

def dung_do_an():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 ---------------------------------------------------------------- bìa
    sl = bia(prs, "ĐỒ ÁN CHUYÊN NGÀNH", "Khoa Công nghệ thông tin & Truyền thông",
             "CMC Restaurant – Hệ thống quản lý nhà hàng,\ngọi món bằng QR và tư vấn món ăn bằng AI",
             "cmcrestaurant.app  ·  order.cmcrestaurant.app  ·  admin.cmcrestaurant.app")
    ghi_chu_noi(sl, """
[~40 giây]
Em chào thầy và các bạn. Nhóm em xin trình bày đồ án chuyên ngành, đề tài CMC
Restaurant — hệ thống quản lý nhà hàng, gọi món bằng mã QR và tư vấn món ăn
bằng AI.

Sản phẩm đang chạy thật trên ba tên miền ghi ở dưới, nên phần trình bày sẽ dùng
ảnh chụp trực tiếp từ hệ thống đang vận hành.
""")

    # 2 ------------------------------------------------------------ lộ trình
    LT = ["Bài toán thực tế và lý do chọn đề tài",
          "Mục tiêu và phạm vi",
          "Sản phẩm và các vai trò sử dụng",
          "Chức năng chính",
          "Kiến trúc và thiết kế hệ thống",
          "Kiểm thử và kết quả đo",
          "Hạn chế và hướng phát triển"]
    sl = lo_trinh(prs, LT)
    ghi_chu_noi(sl, """
[~30 giây]
Phần trình bày gồm bảy phần. Em bắt đầu từ bài toán thực tế, sau đó là mục tiêu
và phạm vi, rồi giới thiệu sản phẩm và các chức năng chính. Phần kiến trúc và
thiết kế là phần kỹ thuật trọng tâm. Cuối cùng là kết quả đo, hạn chế và hướng
phát triển.
""")

    # 3 ------------------------------------------------------------ bài toán
    sl = trang(prs, "Một yêu cầu gọi món đi qua bốn chủ thể", "1 · BÀI TOÁN THỰC TẾ")
    chuoi = ["Khách\nchọn món", "Nhân viên\ntiếp nhận", "Bếp\nchế biến",
             "Quầy\nlập hóa đơn"]
    khe = Cm(1.4)
    w = (RONG - khe * 3) // 4
    for i, ten in enumerate(chuoi):
        x = LE + i * (w + khe)
        o = sl.shapes.add_shape(5, x, Cm(3.7), w, Cm(2.4))
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        tf = o.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for k, d in enumerate(ten.split("\n")):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = d
            r.font.size = Pt(16); r.font.bold = (k == 0)
            r.font.color.rgb = DEN; r.font.name = FONT
        if i < 3:
            mt = sl.shapes.add_shape(13, x + w + Cm(0.2), Cm(4.55), Cm(0.9), Cm(0.8))
            mt.fill.solid(); mt.fill.fore_color.rgb = VANG
            mt.line.fill.background(); mt.shadow.inherit = False
    y_muc(sl, [
        (True, "Mỗi chủ thể chỉ cần một mảnh thông tin, nhưng tất cả phải cùng "
               "tham chiếu đúng bàn, đúng phiên và đúng trạng thái."),
        (True, "Khi thông tin đi bằng lời nói và phiếu giấy, trạng thái không "
               "được ghi lại ở một nơi chung."),
        (False, "Khách chờ nhân viên để hỏi tiến độ · thông tin qua nhiều người "
                "mới tới bếp · quầy cộng tay nhiều lượt gọi khi thanh toán."),
        (True, "Khó khăn thứ hai xuất hiện ngay lúc chọn món: tên món không đủ "
               "để khách biết nguyên liệu hay mức giá."),
    ], y=Cm(7.0), co=16, gian=1.45)
    chu_thich(sl, "Quan sát tại hai nhà hàng ăn tại chỗ quy mô vừa ở Hà Nội, tuần đầu của dự án.")
    ghi_chu_noi(sl, """
[~1 phút 20]
Một yêu cầu gọi món nghe đơn giản nhưng thực tế đi qua bốn chủ thể: khách chọn
món, nhân viên tiếp nhận, bếp chế biến, quầy lập hóa đơn.

Mỗi bên chỉ cần một mảnh thông tin, nhưng cả bốn phải cùng tham chiếu đúng một
bàn, đúng một phiên và đúng một trạng thái.

Khi thông tin đi bằng lời nói và phiếu giấy thì trạng thái ấy không được ghi lại
ở một nơi chung. Hệ quả là khách phải chờ nhân viên mới hỏi được món tới đâu;
thông tin từ bàn phải qua nhiều người mới đến bếp; và lúc thanh toán, quầy phải
cộng tay nhiều lượt gọi.

Còn khó khăn thứ hai xuất hiện ngay lúc chọn món: tên món không đủ để khách biết
nguyên liệu hay mức giá. Câu hỏi thật của khách là "hai người dưới ba trăm nghìn
nên chọn gì", chứ không phải một từ khóa.
""")

    # 4 --------------------------------------------------------- mục tiêu
    sl = trang(prs, "Mục tiêu và phạm vi", "2 · MỤC TIÊU")
    y_muc(sl, [
        (True, "Điểm vào bằng mã QR: nhận diện đúng bàn, mở hoặc tiếp tục đúng phiên."),
        (True, "Luồng gọi món: thực đơn, giỏ phía máy chủ, nhiều lượt, theo dõi trạng thái."),
        (True, "Chuỗi vận hành: bếp cập nhật tiến độ, quầy tổng hợp hóa đơn, quản trị dữ liệu nền."),
        (True, "Một nguồn trạng thái nhất quán, cập nhật gần thời gian thực."),
        (True, "Trợ lý AI có ràng buộc về món, giá và dị nguyên."),
        (True, "Kiểm thử tự động, tích hợp liên tục và môi trường triển khai."),
    ], y=Cm(3.7), w=Cm(19), co=16, gian=1.5)
    x2 = LE + Cm(19.8)
    w2 = RONG - Cm(19.8)
    txt(sl, x2, Cm(3.7), w2, Cm(0.9), "Ngoài phạm vi, có chủ đích", 16, True, NHAN)
    y_muc(sl, [
        (True, "Giao hàng và mang về"),
        (True, "Cổng thanh toán tự động"),
        (True, "Ứng dụng di động cho nhân viên"),
        (True, "Huấn luyện lại mô hình ngôn ngữ"),
    ], y=Cm(4.9), x=x2, w=w2, co=14, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút]
Từ bài toán đó, nhóm đặt sáu mục tiêu, sắp theo đúng thứ tự ưu tiên.

Ba mục tiêu đầu là luồng nghiệp vụ chính: điểm vào bằng QR, luồng gọi món, và
chuỗi vận hành phía nhà hàng.

Mục tiêu thứ tư là điều kiện kỹ thuật cho ba mục tiêu trên: một nguồn trạng thái
nhất quán giữa mọi vai trò.

Mục tiêu thứ năm là trợ lý AI, và em xin nhấn: nó là lớp hỗ trợ, không phải trung
tâm điều khiển.

Mục tiêu thứ sáu là kiểm thử và triển khai — thứ quyết định sản phẩm có kiểm
chứng được hay không.

Bên phải là bốn thứ nhóm cố ý để ngoài phạm vi, ghi rõ để không bị hiểu là thiếu sót.
""")

    # 5 ------------------------------------------------------- sản phẩm
    sl = trang(prs, "Sản phẩm và bốn vai trò sử dụng", "3 · SẢN PHẨM")
    txt(sl, LE, Cm(3.5), RONG, Cm(1.6),
        "Hệ thống web cho nhà hàng ăn tại chỗ: khách quét mã QR trên bàn để xem thực đơn, "
        "hỏi trợ lý AI, gọi món nhiều lượt và theo dõi trạng thái đơn — trong khi bếp, quầy "
        "và quản trị viên làm việc trên cùng một nguồn dữ liệu.",
        17, False, DEN, gian=1.4)
    vai = [("Khách tại bàn", "Quét QR · hỏi AI\ngọi nhiều lượt · theo dõi"),
           ("Bếp", "Hàng đợi món theo trạng thái\ncập nhật tiến độ · hết món"),
           ("Quầy thu ngân", "Một hóa đơn cho cả bàn\nCOD hoặc VietQR · đóng ca"),
           ("Quản trị viên", "Thực đơn · bàn và mã QR\ntài khoản và phân quyền")]
    khe = Cm(0.5)
    w = (RONG - khe * 3) // 4
    for i, (ten, mo) in enumerate(vai):
        x = LE + i * (w + khe)
        o = sl.shapes.add_shape(5, x, Cm(6.0), w, Cm(3.5))
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        txt(sl, x + Cm(0.5), Cm(6.4), w - Cm(1.0), Cm(0.9), ten, 17, True, NHAN)
        txt(sl, x + Cm(0.5), Cm(7.5), w - Cm(1.0), Cm(2), mo, 13, False, DEN, gian=1.3)
    the(sl, [("5", "ứng dụng web"), ("84", "endpoint REST"),
             ("24", "bảng dữ liệu"), ("~97.400", "dòng mã"),
             ("91", "món thực đơn")], y=Cm(10.2), cao=Cm(2.9), co_so=24)
    txt(sl, LE, Cm(13.6), RONG, Cm(2),
        "Sản phẩm chạy thật trên máy chủ, sau HTTPS, với PostgreSQL.\n"
        "Thực đơn hiện tại là dữ liệu mẫu do nhóm dựng để phát triển và đánh giá.",
        15, False, XAM, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút]
Đây là sản phẩm. Một hệ thống web cho nhà hàng ăn tại chỗ, phục vụ bốn vai trò,
mỗi vai trò một giao diện riêng nhưng chung một backend.

Về quy mô: năm ứng dụng web, tám mươi tư endpoint, hai mươi tư bảng dữ liệu,
khoảng chín mươi bảy nghìn dòng mã, thực đơn chín mươi mốt món.

Hai điều em xin nói rõ để thầy đánh giá đúng phạm vi. Thứ nhất, sản phẩm chạy
thật trên máy chủ, sau HTTPS, không phải bản dựng trên máy cá nhân. Thứ hai,
thực đơn hiện tại là dữ liệu mẫu, chưa phải thực đơn vận hành của nhà hàng
thương mại.
""")

    # 6 ----------------------------------------------- chức năng khách + AI
    sl = trang(prs, "Luồng khách tại bàn và trợ lý AI", "4 · CHỨC NĂNG")
    y_muc(sl, [
        (True, "Quét QR mở hoặc tiếp tục đúng phiên; nhiều thiết bị cùng bàn vào chung một phiên."),
        (True, "Giỏ hàng nằm phía máy chủ — đóng tab hay đổi máy vẫn còn."),
        (True, "Gọi món nhiều lượt trong cùng phiên, gộp chung một hóa đơn."),
        (True, "Trạng thái đẩy thời gian thực: Gọi món → Chế biến → Phục vụ → Thanh toán."),
        (True, "Trợ lý AI nhận câu hỏi tiếng Việt theo khẩu vị, ngân sách, dị nguyên."),
        (False, "Thiếu dữ liệu nhãn thì thu hẹp gợi ý và khuyến cáo hỏi nhân viên — không suy đoán."),
    ], y=Cm(3.6), w=Cm(12.0), co=15, gian=1.45)
    luoi_anh(sl, [(RM / "order-scan-2026-07-17.png", "Quét QR vào phiên"),
                  (ANH / "Trangthaibankhach.jpg", "Theo dõi trạng thái"),
                  (ANH / "trolyaitraloi2.jpg", "AI lọc theo dị ứng")],
             Cm(3.6), Cm(12.6), x0=LE + Cm(12.7), w0=RONG - Cm(12.7))
    ghi_chu_noi(sl, """
[~1 phút 40]
Bắt đầu từ phía khách. Khách quét mã QR trên bàn; hệ thống mở phiên mới hoặc
tiếp tục phiên đang mở của đúng bàn đó. Bốn người cùng bàn cùng quét thì vào
chung một phiên chứ không tạo ra bốn phiên.

Giỏ hàng nằm phía máy chủ chứ không trong trình duyệt, nên đóng tab rồi quét lại
vẫn còn giỏ.

Khách gọi món nhiều lượt trong cùng phiên; lượt thứ hai không phải đơn rời mà
gộp chung một hóa đơn.

Ảnh giữa là màn hình theo dõi trạng thái với thanh bốn bước, cập nhật ngay khi
bếp thao tác.

Ảnh bên phải là trợ lý AI. Khách vừa nói dị ứng tôm và danh sách gợi ý thu từ
sáu món xuống ba. Điều đáng chú ý là trợ lý nói thẳng rằng thực đơn chưa ghi
nhận chi tiết thành phần hải sản và đề nghị khách hỏi lại nhân viên — tức là khi
dữ liệu chưa đủ, hệ thống thu hẹp và cảnh báo chứ không đoán.
""")

    # 7 ------------------------------------------------- chức năng vận hành
    sl = trang(prs, "Vận hành nhà hàng và quản trị", "4 · CHỨC NĂNG")
    luoi_anh(sl, [(ANH / "trangbep.jpg", "Bảng bếp — bốn trạng thái"),
                  (ANH / "trangquay.jpg", "Quầy — hóa đơn cả phiên"),
                  (ANH / "quanlythucdon.jpg", "Thực đơn — 91 món"),
                  (ANH / "qrban.jpg", "Sinh mã QR theo bàn")],
             Cm(3.6), Cm(7.6))
    y_muc(sl, [
        (True, "Bếp thao tác vài chạm; bật/tắt hết món có hiệu lực ngay với khách đang xem thực đơn."),
        (True, "Quầy gộp mọi lượt gọi thành một hóa đơn; khuyến mãi áp một lần lúc tất toán."),
        (True, "Quản trị là lớp làm cho vòng đời chạy được: có thực đơn thì AI mới có dữ liệu, "
               "có mã QR thì khách mới vào đúng phiên."),
    ], y=Cm(12.4), co=15, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút 20]
Phía nhà hàng có bốn màn hình.

Bảng bếp chia món theo bốn trạng thái, cập nhật thời gian thực. Bếp thao tác vài
chạm vì tay bận và màn hình đặt xa. Trên bảng này bếp cũng bật tắt được tình
trạng hết món, và thay đổi đó có hiệu lực ngay với khách đang xem thực đơn.

Quầy thu ngân gộp mọi lượt gọi của phiên thành một hóa đơn. Khuyến mãi và tích
điểm áp một lần lúc tất toán chứ không áp theo từng lượt.

Hai ảnh còn lại là phần quản trị: quản lý thực đơn và sinh mã QR theo bàn. Đây
là lớp làm cho toàn bộ vòng đời phía trên chạy được.
""")

    # 8 --------------------------------------------------------- kiến trúc
    sl = trang(prs, "Kiến trúc hệ thống", "5 · THIẾT KẾ")
    anh(sl, SD_DA / "so-do-3.png", LE, Cm(3.5), Cm(19.0), Cm(14.4), vien=False)
    y_muc(sl, [
        (True, "Modular monolith cho nghiệp vụ."),
        (False, "Đơn hàng, thanh toán và phiên bàn luôn đổi cùng nhau nên cần chung một transaction."),
        (True, "Tách riêng duy nhất dịch vụ AI."),
        (False, "Khác ngôn ngữ, ảnh Docker 2,74 GB, nhịp thay đổi khác."),
        (True, "Mọi đường đi của dữ liệu đều qua backend."),
        (False, "Trình duyệt không gọi thẳng dịch vụ AI; dịch vụ AI không chạm cơ sở dữ liệu."),
    ], y=Cm(4.0), x=LE + Cm(19.6), w=RONG - Cm(19.6), co=13, gian=1.3)
    ghi_chu_noi(sl, """
[~1 phút 30]
Về kiến trúc, nhóm chọn modular monolith cho phần nghiệp vụ và chỉ tách riêng
đúng một thứ là dịch vụ AI.

Lý do không dùng microservices: đơn hàng, thanh toán và phiên bàn ở bài toán này
luôn thay đổi cùng nhau và cần nhất quán ngay. Tách ra thì phải dựng saga để mô
phỏng lại thứ mà một transaction vốn cho sẵn.

Nhưng dịch vụ AI được tách, theo tiêu chí khác hẳn: khác ngôn ngữ, ảnh Docker
2,74 GB so với khoảng 200 MB của backend, và nhịp thay đổi khác. Kết luận nhóm
rút ra: tiêu chí tách dịch vụ là khác biệt về vòng đời, không phải sự gọn gàng
của sơ đồ.

Nguyên tắc giữ nhất quán: mọi đường đi của dữ liệu đều qua backend. Nhờ vậy
backend là nơi duy nhất có thẩm quyền về quyền và dữ liệu.
""")

    # 9 ------------------------------------------------ máy trạng thái + AI
    sl = trang(prs, "Hai điểm kỹ thuật cốt lõi", "5 · THIẾT KẾ")
    txt(sl, LE, Cm(3.4), Cm(14), Cm(0.9), "Quét lại thì về đúng bước đang dở",
        17, True, NHAN)
    anh(sl, SD_DA / "so-do-5.png", LE, Cm(4.4), Cm(6.6), Cm(12.4), vien=False)
    y_muc(sl, [
        (True, "Sáu trạng thái được suy ra, không lưu sẵn."),
        (False, "Tính lại từ đơn và hóa đơn mỗi lần quét nên không thể lệch dữ liệu thật."),
        (True, "Không dựa vào thiết bị hay trình duyệt."),
        (False, "Hai thứ đó mất khi khách đổi máy; trạng thái nằm ở phiên bàn trên máy chủ."),
    ], y=Cm(4.6), x=LE + Cm(7.0), w=Cm(7.4), co=13, gian=1.3)

    x2 = LE + Cm(15.2)
    w2 = RONG - Cm(15.2)
    txt(sl, x2, Cm(3.4), w2, Cm(0.9), "Mô hình không được CHỌN món", 17, True, NHAN)
    buoc = [("HIỂU", "đọc tiếng Việt → ràng buộc nhãn"),
            ("CHỌN", "mã TẤT ĐỊNH lọc — mô hình không chạm"),
            ("VIẾT", "diễn đạt trên tập món đã chốt"),
            ("CHẶN", "món hoặc giá ngoài tập → khuôn mẫu")]
    for i, (ten, mo) in enumerate(buoc):
        o = sl.shapes.add_shape(5, x2, Cm(4.5) + i * Cm(1.7), w2, Cm(1.45))
        o.fill.solid()
        o.fill.fore_color.rgb = NEN_DAM if i == 1 else NEN
        o.line.color.rgb = NHAN if i == 1 else KE
        o.shadow.inherit = False
        tf = o.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = ten + "   "
        r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = NHAN if i == 1 else DEN; r.font.name = FONT
        r = p.add_run(); r.text = mo
        r.font.size = Pt(12); r.font.color.rgb = XAM; r.font.name = FONT
    y_muc(sl, [
        (True, "Nếu mô hình chọn món, không cách nào chứng minh nó luôn loại món có tôm cho khách dị ứng tôm."),
        (True, "Khi chọn là phép lọc trên bảng nhãn, câu hỏi thành “bảng nhãn có đúng không” — tra được, kiểm được."),
        (True, "Số đo: lọc theo nhãn đúng 8/8 ca, RAG chỉ 1–2/8."),
    ], y=Cm(11.6), x=x2, w=w2, co=13, gian=1.3)
    ghi_chu_noi(sl, """
[~1 phút 50]
Hai điểm kỹ thuật nhóm cho là cốt lõi.

Bên trái là máy trạng thái phiên bàn. Nó trả lời câu hỏi: khi một người quét mã
của bàn T07, hệ thống lấy gì để quyết định đưa họ tới màn hình nào? Câu trả lời
cố tình không dựa vào thiết bị hay lịch sử trình duyệt, vì hai thứ đó mất khi
khách đổi máy. Nó dựa vào trạng thái của chính phiên bàn trên máy chủ.

Sáu trạng thái này được suy ra chứ không lưu sẵn — tính lại từ danh sách đơn và
trạng thái hóa đơn mỗi lần quét. Cách này chậm hơn chút nhưng loại bỏ hẳn lớp
lỗi mà nhóm đã gặp ở bản đầu: cột trạng thái lưu sẵn bị lệch khỏi dữ liệu thật.

Bên phải là ranh giới quyền của AI. Mô hình được HIỂU câu hỏi và VIẾT câu trả
lời, nhưng bước CHỌN món thì không được chạm vào — đó là mã tất định lọc theo
bảng nhãn.

Vì sao phải tách? Vì nếu để mô hình chọn thì không có cách nào chứng minh nó
luôn loại món có tôm cho khách dị ứng tôm — thử bao nhiêu câu cũng không đủ. Khi
chọn là phép lọc trên bảng nhãn, câu hỏi chuyển thành "bảng nhãn có đúng không",
và điều đó tra được, kiểm được.

Số đo xác nhận: lọc theo nhãn đúng tám trên tám, còn để RAG chọn chỉ đúng một
đến hai trên tám.
""")

    # 10 ------------------------------------------------------- kiểm thử
    sl = trang(prs, "Kiểm thử và kết quả đo", "6 · ĐÁNH GIÁ")
    the(sl, [("118", "test\nfrontend"), ("84", "test\nbackend"),
             ("386", "test mã\ndịch vụ AI"), ("128", "test\nthước đo"),
             ("103/103", "lượt kiểm thử\nđầu-cuối")], y=Cm(3.6), cao=Cm(3.0),
        co_so=24)
    y_muc(sl, [
        (True, "Bốn tầng kiểm thử, cộng một tầng ít ai làm: kiểm chính bộ thước đo."),
        (False, "Bộ dò đưa câu trả lời cố ý vô nghĩa vào thước đo và đòi thước đo cho trượt. "
                "Phát hiện 24 trường hợp chấp nhận sai, đã khắc phục cả 24."),
        (True, "An toàn dị nguyên: không ghi nhận lỗi trên 140 ca + 87 lượt phiên + 8 ca chọn món."),
        (True, "Kiểm thử đầu-cuối dựng toàn bộ hệ thống rồi thao tác như khách thật: 103/103 lượt đạt."),
        (False, "Nó ra đời sau một lỗi mà kiểm thử riêng của cả hai dịch vụ đều xanh nhưng luồng vẫn hỏng."),
    ], y=Cm(7.2), co=16, gian=1.45)
    ghi_chu_noi(sl, """
[~1 phút 20]
Về bảo đảm chất lượng, nhóm có bốn tầng kiểm thử: 118 test frontend, 84 test
backend, 386 test cho mã dịch vụ AI, và 128 test cho chính bộ thước đo. Ngoài ra
có bộ kiểm thử đầu-cuối chạy trên stack thật, đạt 103 trên 103 lượt.

Tầng thứ năm là thứ ít ai làm: kiểm chính bộ thước đo. Phần AI không so được đầu
ra với một giá trị cố định, nó cần hàm chấm điểm — mà hàm chấm điểm cũng là mã,
và mã thì có lỗi. Nhóm viết một bộ dò đưa những câu trả lời cố ý vô nghĩa vào
thước đo và đòi thước đo phải cho trượt. Bộ dò tìm ra 24 trường hợp chấp nhận
sai, và cả 24 đã được khắc phục. Nếu bỏ qua bước này thì mọi con số đều đáng nghi.

Bộ kiểm thử đầu-cuối cũng ra đời từ một sự cố thật: hai dịch vụ trao đổi dữ liệu
theo hai giả định khác nhau về khuôn dạng, khiến luồng hỏng — trong khi kiểm thử
riêng của cả hai đều xanh.
""")

    # 11 -------------------------------------------- hạn chế + hướng phát triển
    sl = trang(prs, "Hạn chế và hướng phát triển", "7 · KẾT LUẬN")
    txt(sl, LE, Cm(3.4), Cm(15.5), Cm(0.9), "Hạn chế của phiên bản hiện tại",
        17, True, NHAN)
    y_muc(sl, [
        (True, "Nhãn dị nguyên mới phủ 44/91 món và chưa được bếp xác nhận."),
        (False, "Fail-closed nên rủi ro nghiêng về thu hẹp gợi ý, nhưng chưa đủ kết luận an toàn y tế."),
        (True, "VietQR mới dừng ở mức sinh mã, quầy xác nhận thủ công."),
        (True, "Độ trễ trợ lý còn cao: p50 8,6 s · p95 13,5 s."),
        (True, "Chưa ước lượng thời gian lên món để báo cho khách."),
        (False, "Cần thời gian chế biến thật từ bếp; nhóm không hiển thị một con số tự đặt."),
        (True, "Khách chưa tự hủy món được — quy tắc đã có, mới thiếu nút phía khách."),
        (True, "Chưa kiểm thử tải, chưa có báo cáo độ phủ; thực đơn là dữ liệu mẫu."),
    ], y=Cm(4.4), w=Cm(15.5), co=14, gian=1.3)
    x2 = LE + Cm(16.3)
    w2 = RONG - Cm(16.3)
    txt(sl, x2, Cm(3.4), w2, Cm(0.9), "Hướng phát triển", 17, True, NHAN)
    gd = [("Ngắn hạn · 1–2 tháng",
           "Hoàn tất nhãn dị nguyên 47 món còn lại và đưa bếp xác nhận · "
           "mở nút hủy món cho khách · kiểm thử tải · báo cáo độ phủ"),
          ("Trung hạn · 3–6 tháng",
           "Ước lượng thời gian lên món theo hướng đo trước hiển thị sau · "
           "giảm p95 dưới 8 giây · webhook ngân hàng cho VietQR"),
          ("Dài hạn",
           "Nhiều chi nhánh · học từ phản hồi khách có vòng kiểm duyệt của người")]
    for i, (ten, mo) in enumerate(gd):
        y = Cm(4.4) + i * Cm(3.4)
        o = sl.shapes.add_shape(5, x2, y, w2, Cm(3.0))
        o.fill.solid()
        o.fill.fore_color.rgb = NEN_DAM if i == 0 else NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        txt(sl, x2 + Cm(0.5), y + Cm(0.25), w2 - Cm(1.0), Cm(0.8), ten, 14, True, NHAN)
        txt(sl, x2 + Cm(0.5), y + Cm(1.1), w2 - Cm(1.0), Cm(1.9), mo, 12, False, DEN, gian=1.28)
    txt(sl, LE, Cm(15.2), RONG, Cm(1.4),
        "Kết quả nhóm coi trọng nhất không phải số tính năng, mà là khả năng truy vết: "
        "từ nhu cầu người dùng tới thiết kế, mã nguồn, kiểm thử và số đo.",
        16, True, DEN, PP_ALIGN.CENTER, gian=1.3)
    txt(sl, LE, Cm(16.9), RONG, Cm(1.2),
        "Xin cảm ơn thầy và các bạn đã lắng nghe.", 19, True, NHAN, PP_ALIGN.CENTER)
    ghi_chu_noi(sl, """
[~1 phút 30]
Cuối cùng là hạn chế và hướng phát triển.

Về hạn chế, quan trọng nhất là nhãn dị nguyên mới phủ 44 trên 91 món và bảng
nhãn chưa được bếp xác nhận. Cơ chế fail-closed khiến rủi ro nghiêng về phía thu
hẹp gợi ý, nhưng em xin nói rõ: điều đó chưa đủ để kết luận hệ thống an toàn về
mặt y tế.

VietQR mới dừng ở mức sinh mã. Độ trễ trợ lý p95 còn 13,5 giây.

Hai hạn chế tiếp theo nằm ở trải nghiệm của khách. Thứ nhất, hệ thống chưa ước
lượng được còn bao lâu thì món lên. Khách thấy được trạng thái đang chuẩn bị hay
đã sẵn sàng, nhưng không biết phải chờ bao lâu. Nhóm có cân nhắc gắn một con số
ước lượng, nhưng đã quyết định không làm, vì muốn ước lượng đáng tin thì cần
thời gian chế biến thật của từng món do bếp cung cấp, mà nhóm chưa có. Hiển thị
một con số tự đặt rồi để khách thấy sai thì hại hơn là không hiển thị gì.

Thứ hai, khách chưa tự hủy món được. Ở đây em xin nói rõ để tránh hiểu nhầm:
quy tắc hủy thì đã có đầy đủ ở tầng nghiệp vụ — món chỉ hủy được khi bếp chưa
làm xong, và cả lượt gọi bị khóa hủy ngay khi một món vào bếp — và màn hình
bếp, phục vụ, quản trị đều đã có nút hủy. Phần còn thiếu chỉ là nút phía khách,
vì endpoint hiện chỉ mở cho tài khoản nhân viên còn khách thì đi bằng thẻ truy
cập theo lượt gọi. Đây là việc rẻ nhất trong danh sách hướng phát triển vì
không phải đổi cơ sở dữ liệu.

Ngoài ra nhóm chưa kiểm thử tải, chưa có báo cáo độ phủ mã nguồn.

Hướng phát triển chia ba giai đoạn, trong đó ngắn hạn tập trung đóng lại chính
những hạn chế vừa nêu.

Kết lại, kết quả nhóm coi trọng nhất không phải số tính năng đã làm, mà là khả
năng truy vết: từ nhu cầu người dùng tới thiết kế, mã nguồn, kiểm thử và số đo.
Mọi con số trong báo cáo đều kèm lệnh chạy lại được.

Phần trình bày của nhóm em đến đây là hết. Em xin cảm ơn và sẵn sàng nhận câu hỏi ạ.
""")

    ra = HERE / "output" / "SLIDE_DO_AN_CHUYEN_NGANH.pptx"
    return _luu(prs, ra)


# ==========================================================================
# BỘ 2 — CÔNG NGHỆ LẬP TRÌNH WEB
# ==========================================================================

def dung_web():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 ---------------------------------------------------------------- bìa
    sl = bia(prs, "BÁO CÁO MÔN HỌC", "Học phần: Công nghệ lập trình Web",
             "CMC Restaurant – Hệ thống quản lý nhà hàng,\ngọi món bằng QR và tư vấn món ăn bằng AI",
             "cmcrestaurant.app  ·  order.cmcrestaurant.app  ·  admin.cmcrestaurant.app")
    ghi_chu_noi(sl, """
[~40 giây]
Em chào thầy và các bạn. Nhóm em xin trình bày báo cáo môn Công nghệ lập trình
Web, đề tài CMC Restaurant — hệ thống quản lý nhà hàng, gọi món bằng mã QR và
tư vấn món ăn bằng AI.

Cách tiếp cận của nhóm là đi từ một sản phẩm thật thay vì trình bày lý thuyết
tách rời: phân tích toàn bộ công nghệ đã dùng để xây dựng hệ thống CMC Restaurant
— khoảng 97.400 dòng mã, đang chạy trên ba tên miền.

Nhờ vậy mọi khẳng định về công nghệ đều kiểm chứng được bằng mã nguồn cụ thể, và
mỗi lựa chọn kỹ thuật đều gắn với một ràng buộc có thật.
""")

    # 2 --------------------------------------------------- bốn thế hệ web
    sl = trang(prs, "Bốn thế hệ kiến trúc ứng dụng web", "1 · TỔNG QUAN")
    bang(sl, ["Thế hệ", "Cách hoạt động", "Hạn chế còn lại"], [
        ["Trang tĩnh", "Máy chủ trả tệp HTML có sẵn", "Không có dữ liệu động"],
        ["Server-rendered", "Máy chủ sinh HTML mỗi lần yêu cầu", "Mỗi thao tác tải lại cả trang"],
        ["SPA + API", "Trình duyệt tự dựng giao diện, dữ liệu qua JSON", "Cần xử lý trạng thái phía client"],
        ["SPA + API + kênh đẩy", "Máy chủ chủ động đẩy dữ liệu xuống", "Phải quản lý trạng thái kết nối"],
    ], y=Cm(3.7), rong_cot=[0.22, 0.44, 0.34], co=15)
    txt(sl, LE, Cm(11.0), RONG, Cm(3),
        "Sản phẩm thuộc thế hệ thứ tư. Lý do gắn với đặc thù bài toán: khi bếp đánh dấu một món "
        "đã xong, màn hình của khách phải đổi ngay — và khách thì không bấm tải lại trang.\n\n"
        "Ba thế hệ đầu đều không đáp ứng được yêu cầu này một cách tự nhiên.",
        17, False, DEN, gian=1.45)
    ghi_chu_noi(sl, """
[~1 phút 10]
Trước hết em xin đặt sản phẩm vào bối cảnh. Kiến trúc web đã đi qua bốn thế hệ,
mỗi thế hệ ra đời để giải một hạn chế của thế hệ trước.

Trang tĩnh chỉ trả về tệp có sẵn. Server-rendered sinh HTML động nhưng mỗi thao
tác phải tải lại cả trang. SPA cộng API giải quyết được điều đó — trình duyệt tự
dựng giao diện, dữ liệu lấy qua JSON.

Nhưng SPA vẫn có một hạn chế: máy chủ không chủ động gửi được dữ liệu. Thế hệ thứ
tư bổ sung một kênh đẩy.

Sản phẩm của nhóm thuộc thế hệ thứ tư, và lý do rất cụ thể: khi bếp đánh dấu một
món đã xong, màn hình của khách phải đổi ngay — mà khách thì không có lý do gì để
bấm tải lại trang.
""")

    # 3 ------------------------------------------------------ sáu tầng
    sl = trang(prs, "Sáu tầng công nghệ của một ứng dụng web", "1 · TỔNG QUAN")
    tang = [("Tầng giao diện", "React 19 · TypeScript · Vite · React Router",
             "5 app · 7 package dùng chung"),
            ("Tầng máy chủ", "ASP.NET Core · Minimal API",
             "84 endpoint · 8 tầng middleware"),
            ("Tầng dữ liệu", "PostgreSQL 16 · Entity Framework Core",
             "24 bảng · 22 migration"),
            ("Kênh thời gian thực", "SignalR", "7 loại sự kiện"),
            ("Lớp bảo mật", "HTTPS · JWT · capability token", "4 vai trò · 4 bộ quét"),
            ("Hạ tầng triển khai", "Docker Compose · GitHub Actions", "9 workflow · 3 tên miền")]
    cao = Cm(2.0)
    for i, (ten, cn, so) in enumerate(tang):
        y = Cm(3.6) + i * (cao + Cm(0.25))
        o = sl.shapes.add_shape(5, LE, y, RONG, cao)
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        txt(sl, LE + Cm(0.6), y + Cm(0.3), Cm(7.0), Cm(1.4), ten, 16, True, NHAN)
        txt(sl, LE + Cm(7.8), y + Cm(0.3), Cm(13.5), Cm(1.4), cn, 15, False, DEN)
        txt(sl, LE + Cm(21.6), y + Cm(0.3), RONG - Cm(22.2), Cm(1.4), so, 14, False, XAM)
    ghi_chu_noi(sl, """
[~1 phút]
Một ứng dụng web đầy đủ gồm sáu nhóm công nghệ, và báo cáo của nhóm đi lần lượt
từng nhóm.

Tầng giao diện: React 19 với TypeScript và Vite, tổ chức thành năm ứng dụng và
bảy thư viện dùng chung.

Tầng máy chủ: ASP.NET Core với Minimal API, 84 endpoint đi qua pipeline tám tầng
middleware.

Tầng dữ liệu: PostgreSQL 16 với Entity Framework Core, 24 bảng qua 22 migration.

Kênh thời gian thực dùng SignalR với bảy loại sự kiện.

Lớp bảo mật gồm HTTPS, JWT và capability token, bốn vai trò và bốn bộ quét tự động.

Cuối cùng là hạ tầng triển khai với Docker Compose và GitHub Actions.

Sáu tầng này chính là bố cục của phần trình bày tiếp theo.
""")

    # 4 ------------------------------------------------------ front-end
    sl = trang(prs, "Tầng giao diện: monorepo năm ứng dụng", "2 · FRONT-END")
    anh(sl, SD_WEB / "so-do-2.png", LE, Cm(3.5), Cm(15.0), Cm(9.5), vien=False)
    y_muc(sl, [
        (True, "Năm ứng dụng, mỗi ứng dụng một bundle riêng."),
        (False, "Khách tại bàn dùng mạng 4G — không có lý do bắt họ tải cả bundle quản trị."),
        (True, "Bảy thư viện dùng chung trong npm workspaces."),
        (False, "shared-types khai báo kiểu một lần; hợp đồng API đổi thì TypeScript báo lỗi "
                "biên dịch ở mọi nơi bị ảnh hưởng, thay vì để lỗi lộ ra lúc chạy."),
        (True, "Vite thay Webpack: dùng ES module gốc, khởi động gần như tức thời."),
    ], y=Cm(3.9), x=LE + Cm(15.6), w=RONG - Cm(15.6), co=14, gian=1.35)
    y_muc(sl, [
        (True, "Định tuyến phía trình duyệt bằng React Router: bấm liên kết không tải lại trang, "
               "History API đổi địa chỉ và chỉ dựng lại phần giao diện tương ứng."),
    ], y=Cm(13.4), co=15, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút 30]
Tầng giao diện là chỗ nhóm có một quyết định kiến trúc đáng nói.

Sản phẩm phục vụ bốn nhóm người dùng có hoàn cảnh khác hẳn nhau: khách dùng điện
thoại qua 4G, bếp dùng màn hình lớn đặt xa, quầy dùng máy tính thao tác nhanh,
quản trị làm việc với bảng nhiều cột.

Một giao diện duy nhất phục vụ cả bốn sẽ hoặc quá nặng cho khách, hoặc quá sơ
sài cho quản trị. Nên nhóm tách thành năm ứng dụng, mỗi ứng dụng một bundle
riêng. Lý do trực tiếp: khách dùng 4G, không có lý do bắt họ tải cả mã của trang
quản trị.

Cái giá phải trả là nguy cơ lặp mã giữa năm ứng dụng. Nhóm giải bằng monorepo với
bảy thư viện dùng chung.

Đáng chú ý nhất là shared-types: kiểu dữ liệu khai báo một lần, cả năm ứng dụng
cùng dùng. Nên khi hợp đồng API đổi, TypeScript báo lỗi biên dịch ở mọi nơi bị
ảnh hưởng — thay vì để lỗi lộ ra lúc chạy, ở đúng màn hình khách đang dùng.
""")

    # 5 -------------------------------------------------- middleware + DI
    sl = trang(prs, "Tầng máy chủ: pipeline xử lý một request", "3 · BACK-END")
    buoc = [("UseForwardedHeaders", "khôi phục IP thật sau reverse proxy"),
            ("UseHsts + HttpsRedirection", "buộc dùng HTTPS"),
            ("UseCors", "kiểm soát nguồn gọi chéo tên miền"),
            ("ApiExceptionHandling", "TỰ VIẾT — chuẩn hóa phản hồi lỗi"),
            ("UseAuthentication", "xác định người gọi là ai"),
            ("UseAuthorization", "quyết định người gọi được làm gì"),
            ("Endpoint", "84 endpoint nghiệp vụ")]
    cao = Cm(1.55)
    for i, (ten, mo) in enumerate(buoc):
        y = Cm(3.7) + i * (cao + Cm(0.2))
        noi_bat = (i == 3)
        o = sl.shapes.add_shape(5, LE, y, Cm(17.0), cao)
        o.fill.solid()
        o.fill.fore_color.rgb = NEN_DAM if noi_bat else NEN
        o.line.color.rgb = NHAN if noi_bat else KE
        o.shadow.inherit = False
        tf = o.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"  {ten}   "
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = NHAN if noi_bat else DEN; r.font.name = FONT
        r = p.add_run(); r.text = mo
        r.font.size = Pt(12); r.font.color.rgb = XAM; r.font.name = FONT
    y_muc(sl, [
        (True, "Thứ tự là một phần của thiết kế."),
        (False, "UseAuthentication phải đứng trước UseAuthorization — không thể quyết định "
                "một người được làm gì khi chưa biết họ là ai."),
        (True, "Middleware tự viết trả mã tham chiếu 8 ký tự."),
        (False, "Người vận hành tra được log bằng mã đó, còn kẻ tấn công không đọc được "
                "cấu trúc nội bộ của hệ thống."),
        (True, "Dependency Injection đăng ký theo mô-đun nghiệp vụ."),
        (False, "Nhờ vậy tệp khởi động chỉ 273 dòng dù hệ thống có 84 endpoint."),
    ], y=Cm(3.9), x=LE + Cm(17.8), w=RONG - Cm(17.8), co=13, gian=1.3)
    ghi_chu_noi(sl, """
[~1 phút 40]
Middleware pipeline là khái niệm trung tâm của ASP.NET Core. Mỗi yêu cầu HTTP đi
qua một chuỗi thành phần; mỗi thành phần xử lý rồi chuyển tiếp cho thành phần kế
tiếp.

Đây là tám tầng thật trong dự án. Em xin nhấn hai điểm.

Thứ nhất, thứ tự là một phần của thiết kế, không phải ngẫu nhiên.
UseAuthentication phải đứng trước UseAuthorization, vì không thể quyết định một
người được làm gì khi chưa biết họ là ai.

Thứ hai, tầng được tô đậm là middleware nhóm tự viết. Nó bắt mọi ngoại lệ chưa
xử lý, ghi log đầy đủ, nhưng chỉ trả ra ngoài một mã tham chiếu tám ký tự. Người
vận hành tra được log bằng mã đó, còn người dùng — kể cả kẻ tấn công — không đọc
được cấu trúc nội bộ.

Về Dependency Injection, dự án đăng ký dịch vụ theo mô-đun nghiệp vụ thay vì liệt
kê phẳng. Nhờ vậy tệp khởi động chỉ 273 dòng dù hệ thống có 84 endpoint.
""")

    # 6 ---------------------------------------------------------- REST API
    sl = trang(prs, "Thiết kế REST API", "3 · BACK-END")
    bang(sl, ["Động từ", "Đường dẫn ví dụ", "Ý nghĩa", "Mã trả về"], [
        ["GET", "/api/menu", "Lấy danh sách thực đơn", "200"],
        ["GET", "/api/orders/{id}", "Lấy một đơn cụ thể", "200 · 404"],
        ["POST", "/api/tables/scan", "Mở hoặc tiếp tục phiên bàn", "200"],
        ["POST", "/api/orders", "Tạo một lượt đơn mới", "201"],
        ["DELETE", "/api/cart/items/{id}", "Xóa một dòng giỏ hàng", "204"],
    ], y=Cm(3.6), rong_cot=[0.12, 0.30, 0.42, 0.16], co=14)
    y_muc(sl, [
        (True, "Minimal API thay vì MVC Controller — 0 tệp kế thừa ControllerBase."),
        (False, "Tầng máy chủ không kết xuất giao diện, nên toàn bộ bộ máy View, Razor và "
                "HTML Helper của MVC không dùng đến. Chọn Minimal API là bỏ đi phần không dùng."),
        (True, "Khóa idempotency chống tạo đơn trùng khi mạng chập chờn."),
        (False, "Khách bấm gửi món, mạng ngắt, ứng dụng gửi lại — bếp nhận hai đơn giống nhau. "
                "Mỗi yêu cầu mang một khóa duy nhất; gửi lại cùng khóa thì trả kết quả cũ."),
        (True, "Mã lỗi nghiệp vụ chuẩn hóa: CONFLICT_STALE, PROMOTION_NOT_FOUND…"),
    ], y=Cm(10.0), co=15, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút 30]
Về thiết kế API, dự án theo REST: mỗi đường dẫn đại diện một tài nguyên, động từ
HTTP cho biết thao tác.

Điểm đáng nói thứ nhất: dự án dùng Minimal API, không có tệp nào kế thừa
ControllerBase. Lý do là tầng máy chủ không kết xuất giao diện — nó chỉ trả JSON
cho năm ứng dụng React. Toàn bộ bộ máy View, Razor, HTML Helper của MVC vì vậy
không dùng đến.

Em xin nói rõ để tránh hiểu nhầm: MVC Controller không hề lỗi thời. Nếu ứng dụng
cần kết xuất HTML phía máy chủ hoặc tối ưu SEO thì MVC với Razor là lựa chọn tự
nhiên hơn.

Điểm thứ hai là khóa idempotency. Mạng di động chập chờn dẫn tới vấn đề thực tế:
khách bấm gửi món, mạng ngắt, ứng dụng gửi lại, và bếp nhận hai đơn giống nhau.
Mỗi yêu cầu mang một khóa duy nhất; máy chủ nhận lại đúng khóa đó thì trả kết quả
cũ thay vì tạo bản ghi mới. Đây là kỹ thuật bắt buộc với ứng dụng có liên quan
tới tiền.
""")

    # 7 --------------------------------------------------------- dữ liệu
    sl = trang(prs, "Tầng dữ liệu và ba bất biến", "4 · DỮ LIỆU")
    anh(sl, SD_WEB / "so-do-3.png", LE, Cm(3.5), Cm(11.0), Cm(14.2), vien=False)
    x2 = LE + Cm(11.6)
    w2 = RONG - Cm(11.6)
    txt(sl, x2, Cm(3.5), w2, Cm(0.9),
        "Ba quy tắc đặt ở tầng cơ sở dữ liệu, không ở tầng ứng dụng", 16, True, NHAN)
    bd = [("Một bàn chỉ có một phiên đang mở", "Unique index có điều kiện",
           "Bốn người quét cùng lúc thì bốn tiến trình cùng đọc “chưa có phiên” rồi cùng tạo"),
          ("Mã đơn không trùng khi nhiều máy chủ", "PostgreSQL sequence",
           "Sinh mã phía ứng dụng thì hai máy chủ có thể sinh cùng một số"),
          ("Hai người sửa cùng đơn không ghi đè", "Cột hệ thống xmin",
           "Người sửa sau sẽ âm thầm xóa thay đổi của người sửa trước")]
    for i, (bt, cc, vs) in enumerate(bd):
        y = Cm(4.8) + i * Cm(3.6)
        o = sl.shapes.add_shape(5, x2, y, w2, Cm(3.2))
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        txt(sl, x2 + Cm(0.5), y + Cm(0.25), w2 - Cm(1.0), Cm(0.8), bt, 14, True, DEN)
        txt(sl, x2 + Cm(0.5), y + Cm(1.05), w2 - Cm(1.0), Cm(0.7), cc, 13, True, NHAN)
        txt(sl, x2 + Cm(0.5), y + Cm(1.85), w2 - Cm(1.0), Cm(1.2), vs, 12, False, XAM, gian=1.25)
    txt(sl, x2, Cm(15.8), w2, Cm(1.4),
        "Nguyên tắc: tầng ứng dụng có lỗi, tầng cơ sở dữ liệu thì không.",
        15, True, DEN, gian=1.3)
    ghi_chu_noi(sl, """
[~1 phút 40]
Tầng dữ liệu dùng PostgreSQL 16 với Entity Framework Core: 24 bảng hình thành
qua 22 migration.

Lựa chọn PostgreSQL không phải mặc định mà xuất phát từ ba yêu cầu cụ thể — ba
bất biến bên phải.

Thứ nhất, một bàn chỉ được có một phiên đang mở. Nếu kiểm ở tầng ứng dụng thì
khi bốn người quét cùng lúc, bốn tiến trình cùng đọc thấy "chưa có phiên" rồi
cùng tạo. Chỉ unique index có điều kiện ở cơ sở dữ liệu mới chặn được.

Thứ hai, mã đơn không được trùng khi chạy nhiều máy chủ. Sinh mã phía ứng dụng
thì hai máy chủ có thể sinh cùng một số, nên dùng sequence của PostgreSQL.

Thứ ba, hai người sửa cùng một đơn không được ghi đè nhau. Dự án dùng cột hệ
thống xmin của PostgreSQL làm concurrency token: khi lưu, EF Core thêm điều kiện
kiểm tra phiên bản. Nếu dòng đã bị người khác sửa thì câu lệnh không khớp và hệ
thống trả lỗi thay vì làm mất dữ liệu.

Nguyên tắc chung rút ra: tầng ứng dụng có lỗi, tầng cơ sở dữ liệu thì không.
Ràng buộc nào quan trọng tới mức không được phép sai thì nên đặt ở nơi thấp nhất
có thể.
""")

    # 8 --------------------------------------------------- thời gian thực
    sl = trang(prs, "Giao tiếp thời gian thực", "5 · REALTIME")
    bang(sl, ["Cơ chế", "Cách hoạt động", "Hạn chế"], [
        ["Polling", "Client hỏi lại theo chu kỳ", "Trễ bằng chu kỳ hỏi; tốn băng thông"],
        ["Server-Sent Events", "Máy chủ giữ kết nối và đẩy sự kiện xuống", "Một chiều"],
        ["WebSocket", "Kênh song công trên một kết nối TCP", "Có thể bị proxy chặn"],
    ], y=Cm(3.6), rong_cot=[0.22, 0.44, 0.34], co=14)
    y_muc(sl, [
        (True, "SignalR trừu tượng hóa cả ba: ưu tiên WebSocket, tự lùi về SSE rồi long-polling."),
        (False, "Quan trọng trong thực tế — mạng wifi nhà hàng thường đi qua proxy, "
                "và không biết trước proxy có cho WebSocket đi qua hay không."),
        (True, "Bảy loại sự kiện, gửi theo nhóm chứ không phát rộng."),
        (False, "order.created · order.statusChanged · order.itemStatusChanged · cart.updated · "
                "payment.requested · assistance.requested · menu.availabilityChanged"),
        (True, "Trạng thái khách thấy và trạng thái bếp thao tác là CÙNG MỘT bản ghi, "
               "hiển thị theo hai vai trò."),
    ], y=Cm(8.6), co=15, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút 30]
Mô hình HTTP truyền thống là client hỏi, server trả lời — máy chủ không chủ động
gửi được. Với bài toán này thì đó là hạn chế chí mạng: khi bếp đánh dấu món đã
xong, màn hình khách phải đổi ngay.

Có ba cơ chế. Polling thì đơn giản nhưng trễ và tốn băng thông. Server-Sent
Events nhẹ nhưng một chiều. WebSocket độ trễ thấp và hai chiều nhưng có thể bị
proxy chặn.

Dự án dùng SignalR, thư viện trừu tượng hóa cả ba. Điểm mạnh quyết định là cơ
chế lùi tự động: ưu tiên WebSocket, nếu bị chặn thì chuyển sang SSE, rồi
long-polling — mà mã ứng dụng không phải thay đổi gì. Điều này quan trọng vì
mạng wifi nhà hàng thường đi qua proxy và không biết trước proxy có cho WebSocket
đi qua hay không.

Hệ thống có bảy loại sự kiện, và chúng được gửi theo nhóm chứ không phát rộng —
sự kiện của bàn T07 chỉ tới thiết bị đang mở phiên bàn T07.

Điểm cần nhấn: trạng thái khách thấy và trạng thái bếp thao tác không phải hai
bản sao cần đồng bộ. Chúng là cùng một bản ghi, hiển thị theo hai vai trò.
""")

    # 9 ------------------------------------------------- nhiều dịch vụ + SSE
    sl = trang(prs, "Kiến trúc nhiều dịch vụ và truyền luồng", "6 · DỊCH VỤ")
    y_muc(sl, [
        (True, "Tiêu chí tách dịch vụ là khác biệt vòng đời, không phải sơ đồ đẹp."),
        (False, "Dịch vụ AI khác ngôn ngữ, ảnh Docker 2,74 GB so với 200 MB, nhịp thay đổi khác. "
                "Gộp chung thì sửa một endpoint thực đơn cũng phải đóng gói lại cả thư viện xử lý ngôn ngữ."),
        (True, "Ba ràng buộc khi hai dịch vụ nói chuyện với nhau."),
        (False, "Trình duyệt không gọi thẳng · có token xác thực nội bộ · "
                "dịch vụ AI không có kết nối tới cơ sở dữ liệu."),
        (True, "Truyền luồng bằng Server-Sent Events cho câu trả lời dài."),
        (False, "Chữ hiện dần thay vì màn hình trống vài giây. Từng có sự cố thiếu dòng "
                "event: khiến luồng hỏng — trong khi test riêng của cả hai dịch vụ đều xanh."),
    ], y=Cm(3.6), w=Cm(20), co=15, gian=1.4)
    x2 = LE + Cm(20.8)
    w2 = RONG - Cm(20.8)
    txt(sl, x2, Cm(3.6), w2, Cm(0.9), "Ngân sách thời gian chờ", 16, True, NHAN)
    ns = [("Trình duyệt", "60 s"), ("Máy chủ nghiệp vụ", "50 s"), ("Dịch vụ AI", "30 s")]
    for i, (ten, so) in enumerate(ns):
        y = Cm(4.8) + i * Cm(2.3)
        o = sl.shapes.add_shape(5, x2 + Cm(i * 0.4), y, w2 - Cm(i * 0.4), Cm(1.9))
        o.fill.solid(); o.fill.fore_color.rgb = NEN
        o.line.color.rgb = KE; o.shadow.inherit = False
        txt(sl, x2 + Cm(i * 0.4 + 0.4), y + Cm(0.2), w2 - Cm(1.2), Cm(0.7), ten, 13, True, DEN)
        txt(sl, x2 + Cm(i * 0.4 + 0.4), y + Cm(0.95), w2 - Cm(1.2), Cm(0.7), so, 15, True, NHAN)
    txt(sl, x2, Cm(11.8), w2, Cm(3),
        "Giảm dần từ ngoài vào trong.\n\nNếu tầng trong chờ lâu hơn tầng ngoài, "
        "tầng ngoài bỏ cuộc trước và trả lỗi trong khi tầng trong vẫn đang xử lý.",
        13, False, XAM, gian=1.3)
    ghi_chu_noi(sl, """
[~1 phút 30]
Khi hệ thống có nhiều dịch vụ, câu hỏi kỹ thuật là: khi nào thì tách?

Dự án giữ toàn bộ nghiệp vụ trong một tiến trình và chỉ tách đúng một thứ là dịch
vụ tư vấn. Tiêu chí không phải sự gọn gàng của sơ đồ mà là khác biệt về vòng đời:
khác ngôn ngữ, ảnh Docker 2,74 GB so với 200 MB, nhịp thay đổi khác. Gộp chung
thì mỗi lần sửa một endpoint thực đơn cũng phải đóng gói lại toàn bộ thư viện xử
lý ngôn ngữ.

Khi hai dịch vụ nói chuyện, nhóm đặt ba ràng buộc: trình duyệt không bao giờ gọi
thẳng dịch vụ thứ hai; mọi lời gọi nội bộ có token xác thực riêng; và dịch vụ đó
không có kết nối tới cơ sở dữ liệu.

Về truyền luồng: câu trả lời mất vài giây để sinh. Nếu chờ xong mới trả thì màn
hình trống, dễ tưởng treo. Server-Sent Events cho phép gửi từng phần ngay khi có.

Bên phải là ngân sách thời gian chờ, giảm dần từ ngoài vào trong. Nếu đặt ngược
lại — tầng trong chờ lâu hơn tầng ngoài — thì tầng ngoài bỏ cuộc trước, tài
nguyên bị chiếm vô ích và thông báo lỗi không phản ánh đúng nguyên nhân.
""")

    # 10 --------------------------------------------------------- bảo mật
    sl = trang(prs, "Bảo mật ứng dụng web", "7 · BẢO MẬT")
    bang(sl, ["Lớp", "Cơ chế"], [
        ["Kênh truyền", "HTTPS bắt buộc · HSTS · chuyển hướng HTTP sang HTTPS"],
        ["Gọi chéo tên miền", "CORS khai báo danh sách nguồn được phép"],
        ["Xác thực nhân viên", "JWT Bearer — stateless, mở rộng ngang dễ"],
        ["Xác thực khách", "Capability token cấp theo từng lần quét mã QR"],
        ["Lưu mật khẩu", "PBKDF2-HMAC-SHA256 có salt · khóa tài khoản khi sai nhiều lần"],
        ["Quét lỗ hổng", "CodeQL · gitleaks · Trivy · dependency-review"],
    ], y=Cm(3.6), rong_cot=[0.26, 0.74], co=14)
    y_muc(sl, [
        (True, "Mã định danh không phải chứng chỉ ủy quyền."),
        (False, "Biết sessionId không đồng nghĩa với có quyền thao tác. Mã trong QR được in và "
                "dán công khai tại bàn, nên phải đổi lấy capability token cấp riêng cho từng lần quét."),
        (True, "Không tin client: vai trò trong React chỉ để hiển thị cho gọn mắt."),
        (False, "Ẩn một nút không ngăn được ai đó gọi thẳng API bằng công cụ khác."),
        (True, "Quét tự động đã bắt một lỗ hổng thật mà ba người đọc mã đều bỏ sót."),
    ], y=Cm(10.6), co=15, gian=1.4)
    ghi_chu_noi(sl, """
[~1 phút 30]
Bảo mật không nằm ở một cơ chế duy nhất mà ở nhiều lớp bổ trợ nhau.

Em xin nhấn hai điểm.

Thứ nhất là bài toán riêng của sản phẩm: khách tại bàn không có tài khoản, họ chỉ
quét mã QR. Vậy dựa vào đâu để cho phép họ thao tác? Cách ngây thơ là dùng chính
sessionId làm chứng chỉ — ai biết mã phiên thì được thao tác. Cách này sai, vì mã
phiên xuất hiện trên thanh địa chỉ.

Nhóm tách bạch hai khái niệm: mã định danh không phải chứng chỉ ủy quyền. Mã trong
QR được in và dán công khai tại bàn, nên nó chỉ dùng để đổi lấy một capability
token cấp riêng cho từng lần quét.

Thứ hai là nguyên tắc không tin client. Vai trò trong ứng dụng React chỉ dùng để
quyết định hiển thị gì cho gọn mắt, không phải cơ chế phân quyền. Ẩn một nút
không ngăn được ai đó gọi thẳng API bằng công cụ khác.

Về quét tự động: CodeQL đã báo một lỗ hổng thật tại ba vị trí mà không thành viên
nào nhận ra khi đọc lại mã. Bài học: con người rà soát theo ý định, còn công cụ
rà soát theo luồng dữ liệu.
""")

    # 11 ----------------------------------------------------------- CI/CD
    sl = trang(prs, "CI/CD và triển khai", "8 · TRIỂN KHAI")
    anh(sl, SD_WEB / "so-do-4.png", LE, Cm(3.5), Cm(14.5), Cm(13.6), vien=False)
    y_muc(sl, [
        (True, "CI — mỗi thay đổi được kiểm tự động trong vài phút."),
        (False, "Năm job chạy song song: build giao diện, kiểm thử máy chủ, kiểm thử dữ liệu, "
                "kiểm thử đầu-cuối, kiểm tra cấu hình."),
        (True, "CD — triển khai tự động, không thao tác tay."),
        (False, "Loại bỏ lớp lỗi do triển khai thủ công: quên một bước, chép nhầm cấu hình."),
        (True, "CI phải là cổng chặn, không phải bảng thông báo."),
        (False, "Năm job là status check bắt buộc, danh sách ngoại lệ để rỗng — "
                "kể cả chủ repository cũng không bỏ qua được."),
        (True, "Kiểm tra sức khỏe hỏng thì tự động quay về phiên bản trước."),
        (False, "Không chờ người bấm, vì sự cố có thể xảy ra lúc không ai trực."),
    ], y=Cm(3.7), x=LE + Cm(15.1), w=RONG - Cm(15.1), co=13, gian=1.28)
    ghi_chu_noi(sl, """
[~1 phút 30]
Cuối cùng là triển khai. Em xin phân biệt rõ hai khái niệm hay bị dùng lẫn.

CI — tích hợp liên tục — là mỗi thay đổi mã nguồn được tự động kiểm tra ngay khi
đưa lên: biên dịch được không, test có đạt không, cấu hình có hợp lệ không. Mục
đích là phát hiện lỗi trong vài phút thay vì vài ngày.

CD — triển khai liên tục — là sau khi CI đạt thì phiên bản mới tự động lên môi
trường chạy. Điều này loại bỏ lớp lỗi do triển khai thủ công: quên một bước, chép
nhầm tệp cấu hình, hoặc mỗi lần làm một kiểu.

Dự án có chín workflow và 2.468 lần chạy.

Điểm em muốn nhấn: một pipeline CI chỉ có giá trị khi nó ngăn được mã lỗi vào
nhánh chính. Nếu nó chỉ báo đỏ mà vẫn merge được thì nó là bảng thông báo, không
phải cổng kiểm soát. Nên năm job được khai báo là status check bắt buộc, và danh
sách ngoại lệ để rỗng — kể cả chủ repository cũng không bỏ qua được.

Nếu kiểm tra sức khỏe sau triển khai thất bại, hệ thống tự động quay về phiên bản
trước, không chờ người bấm.
""")

    # 12 --------------------------------------------------------- tổng kết
    sl = trang(prs, "Tổng kết")
    y_muc(sl, [
        (True, "Mỗi lựa chọn công nghệ đều xuất phát từ một ràng buộc có thật."),
        (False, "Tách năm ứng dụng vì khách dùng 4G · chọn PostgreSQL vì cần ba tính năng cụ thể · "
                "chọn SignalR vì mạng nhà hàng có thể chặn WebSocket."),
        (True, "Ngược lại, không dùng Kubernetes và không dùng cơ sở dữ liệu vector — "
               "vì không có ràng buộc nào đòi hỏi chúng."),
        (True, "Tầng thấp nhất có thể là nơi tốt nhất để đặt ràng buộc quan trọng."),
        (False, "Tầng ứng dụng có lỗi; tầng cơ sở dữ liệu thì không."),
        (True, "Không tin dữ liệu từ phía client — mọi kiểm tra quan trọng phải lặp lại ở máy chủ."),
    ], y=Cm(3.8), co=17, gian=1.5)
    txt(sl, LE, Cm(14.6), RONG, Cm(1.4),
        "Xin cảm ơn thầy và các bạn đã lắng nghe.", 22, True, NHAN, PP_ALIGN.CENTER)
    txt(sl, LE, Cm(16.2), RONG, Cm(1),
        "github.com/Anpham120/restaurant-qr-ai-ordering", 14, False, XAM, PP_ALIGN.CENTER)
    ghi_chu_noi(sl, """
[~1 phút]
Ba bài học nhóm rút ra sau khi xây dựng hệ thống này.

Thứ nhất, mỗi lựa chọn công nghệ nên xuất phát từ một ràng buộc có thật. Nhóm
tách năm ứng dụng vì khách dùng 4G; chọn PostgreSQL vì cần đúng ba tính năng cụ
thể; chọn SignalR vì mạng nhà hàng có thể chặn WebSocket. Ngược lại, nhóm không
dùng Kubernetes, không dùng cơ sở dữ liệu vector chuyên dụng — vì không có ràng
buộc nào đòi hỏi chúng.

Thứ hai, tầng thấp nhất có thể là nơi tốt nhất để đặt ràng buộc quan trọng. Tầng
ứng dụng có lỗi; tầng cơ sở dữ liệu thì không.

Thứ ba, không tin dữ liệu từ phía client. Giao diện có thể bị sửa bằng công cụ
phát triển của trình duyệt, nên mọi kiểm tra quan trọng phải lặp lại ở máy chủ.

Phần trình bày của nhóm em đến đây là hết. Em xin cảm ơn và sẵn sàng nhận câu hỏi ạ.
""")

    ra = HERE / "output" / "SLIDE_CONG_NGHE_LAP_TRINH_WEB.pptx"
    return _luu(prs, ra)


def _luu(prs, ra):
    ra.parent.mkdir(parents=True, exist_ok=True)
    dich = ra
    try:
        prs.save(dich)
    except PermissionError:
        from datetime import datetime
        dich = ra.with_name(f"{ra.stem}_{datetime.now():%H%M%S}.pptx")
        prs.save(dich)
        print(f"  CẢNH BÁO: {ra.name} đang mở trong PowerPoint.")
        print(f"            Đã lưu sang {dich.name}.")
    print(f"  {dich.name}: {len(prs.slides._sldIdLst)} slide")
    return dich


if __name__ == "__main__":
    print("Đang dựng slide...")
    dung_do_an()
    dung_web()
    print("Xong. Kịch bản nói nằm ở phần ghi chú của từng slide.")
