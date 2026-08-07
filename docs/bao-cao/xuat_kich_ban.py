"""Sinh kịch bản thuyết trình từ chính phần ghi chú của tệp .pptx.

Ghi chú slide là nguồn duy nhất: sửa kịch bản thì sửa trong xuat_slide_pptx.py
rồi dựng lại slide và chạy script này, hai bản không thể lệch nhau.

Cách dùng:
    python xuat_kich_ban.py            -> sinh cho cả hai bộ slide
"""
import re
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
RA = HERE / "output"

BO = [
    ("SLIDE_DO_AN_CHUYEN_NGANH.pptx", "KICH_BAN_DO_AN.md",
     "Kịch bản thuyết trình — Đồ án chuyên ngành"),
    ("SLIDE_CONG_NGHE_LAP_TRINH_WEB.pptx", "KICH_BAN_WEB.md",
     "Kịch bản thuyết trình — Công nghệ lập trình Web"),
]


def tieu_de_slide(sl, thu_tu: int) -> str:
    """Lấy dòng chữ lớn nhất trên slide làm tên slide."""
    ung_vien = []
    for sh in sl.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        co = max((r.font.size.pt for p in sh.text_frame.paragraphs
                  for r in p.runs if r.font.size), default=0)
        ung_vien.append((co, re.sub(r"\s+", " ", sh.text_frame.text).strip()))
    if not ung_vien:
        return f"Slide {thu_tu}"
    chu = max(ung_vien)[1]
    return chu if len(chu) <= 72 else chu[:69].rstrip(" ,·") + "…"


def dung(ten_pptx: str, ten_md: str, tieu_de: str):
    prs = Presentation(str(RA / ten_pptx))
    phan, tong_tu = [], 0

    for i, sl in enumerate(prs.slides, 1):
        ghi_chu = ""
        if sl.has_notes_slide:
            ghi_chu = sl.notes_slide.notes_text_frame.text.strip()
        tong_tu += len(ghi_chu.split())
        phan.append(f"## Slide {i} — {tieu_de_slide(sl, i)}\n\n{ghi_chu or '_(chưa có ghi chú)_'}\n")

    # ~130 từ/phút là tốc độ nói tiếng Việt thong thả khi thuyết trình.
    phut = tong_tu / 130
    dau = [
        f"# {tieu_de}\n",
        f"Bản này sinh tự động từ phần **ghi chú** của `{ten_pptx}` "
        f"(`xuat_kich_ban.py`), nên luôn khớp với slide đang dùng. "
        f"Trong lúc trình bày, mở **Presenter View** của PowerPoint là thấy "
        f"đúng nội dung này.\n",
        f"**{len(prs.slides._sldIdLst)} slide · {tong_tu} từ · "
        f"khoảng {phut:.0f}–{phut * 1.25:.0f} phút** tùy tốc độ nói.\n",
        "---\n",
    ]
    (HERE / ten_md).write_text("\n".join(dau + phan), encoding="utf-8")
    print(f"{ten_md}: {len(prs.slides._sldIdLst)} slide · {tong_tu} từ · "
          f"~{phut:.0f}–{phut * 1.25:.0f} phút")


for a, b, c in BO:
    dung(a, b, c)
